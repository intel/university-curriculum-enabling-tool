#!/usr/bin/env python3
import json
import sys
import re
import math
import io
import base64
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.fill import FillFormat

# Enhanced color palette with complementary colors
COLORS = {
    # Primary colors
    "primary": RGBColor(37, 99, 235),     # Blue
    "primary_light": RGBColor(96, 165, 250),  # Light blue
    "primary_dark": RGBColor(30, 64, 175),   # Dark blue
    
    # Secondary colors
    "secondary": RGBColor(79, 70, 229),   # Indigo
    "secondary_light": RGBColor(139, 92, 246),  # Light indigo/violet
    "secondary_dark": RGBColor(67, 56, 202),    # Dark indigo
    
    # Accent colors
    "accent1": RGBColor(139, 92, 246),     # Violet
    "accent2": RGBColor(16, 185, 129),     # Emerald
    "accent3": RGBColor(245, 158, 11),     # Amber
    "accent4": RGBColor(239, 68, 68),      # Red
    
    # Neutral colors
    "light": RGBColor(243, 244, 246),     # Light gray
    "light_alt": RGBColor(249, 250, 251), # Off-white
    "dark": RGBColor(31, 41, 55),         # Dark gray
    "dark_alt": RGBColor(17, 24, 39),     # Near black
    
    # Text colors
    "text": RGBColor(17, 24, 39),         # Near black
    "text_light": RGBColor(255, 255, 255), # White
    "text_muted": RGBColor(107, 114, 128), # Medium gray
    
    # Status colors
    "success": RGBColor(16, 185, 129),    # Green
    "warning": RGBColor(245, 158, 11),    # Amber
    "error": RGBColor(239, 68, 68),       # Red
    "info": RGBColor(59, 130, 246),       # Blue
    
    # Background colors
    "background": RGBColor(255, 255, 255), # White
    "background_alt": RGBColor(249, 250, 251), # Off-white
    
    # Theme colors
    "royal_blue": RGBColor(65, 105, 225),  # Royal blue
    "medium_purple": RGBColor(147, 112, 219), # Medium purple
    "dark_blue": RGBColor(26, 43, 60),     # Dark blue/black
    "teal": RGBColor(20, 184, 166),        # Teal
    "emerald": RGBColor(16, 185, 129),     # Emerald
    "gradient_start": RGBColor(65, 105, 225), # Royal blue
    "gradient_end": RGBColor(147, 112, 219),  # Medium purple
    
    # Activity slide colors
    "activity_purple": RGBColor(139, 92, 246),  # Purple for activity badge
    "activity_blue": RGBColor(37, 99, 235),     # Blue for header
    "activity_green": RGBColor(16, 185, 129),   # Green for materials accent
    "activity_orange": RGBColor(249, 115, 22),  # Orange for corner triangle
}

# Constants for slide layout
SLIDE_WIDTH = 10  # inches
SLIDE_HEIGHT = 5.625  # inches
FOOTER_Y = 5.3  # Y position for footer elements
CONTENT_START_Y = 1.5  # Starting Y position for content
AVAILABLE_CONTENT_HEIGHT = FOOTER_Y - CONTENT_START_Y  # Available height for content
MAIN_BULLET_INDENT = 0.5  # Left indent for main bullets
SUB_BULLET_INDENT = 1.0   # Left indent for sub-bullets
SUB_SUB_BULLET_INDENT = 1.5  # Left indent for sub-sub-bullets

# Visual theme settings
THEME = {
    "use_gradients": True,
    "corner_accent": True,
    "slide_border": False,
    "content_box_shadow": True,
    "modern_bullets": True,
    "footer_style": "modern",  # "modern" or "classic"
}

# Bullet point markers for detection
BULLET_MARKERS = ['•', '*', '-', '○', '◦', '▪', '▫', '◆', '◇', '►', '▻', '▶', '▷']
SUB_BULLET_MARKERS = ['-', '○', '◦', '▪', '▫']

def clean_slide_title(title):
    """
    Remove slide numbers, colons, and clean up the title.
    This approach handles patterns like 'Slide X:', 'Slide X -', or any text before a colon.
    """
    # Check for colon in the title
    if ':' in title:
        # Split by colon and take everything after it
        title = title.split(':', 1)[1].strip()
        return title
        
    # If no colon, check for the "Slide X" pattern
    title_parts = title.split()
    
    # Check if the title starts with "Slide" followed by a number
    if len(title_parts) > 1 and title_parts[0].lower() == "slide" and title_parts[1].replace(":", "").isdigit():
        # Remove the first two parts ("Slide" and the number)
        title_parts = title_parts[2:]
        # Join the remaining parts back into a string
        return " ".join(title_parts).strip()
    
    return title.strip()

def clean_activity_title(title):
    """
    Remove slide numbers, colons, and clean up the title.
    This approach handles patterns like 'Slide X:', 'Slide X -', or any text before a colon.
    """
    # Check for colon in the title
    if ':' in title:
        # Split by colon and take everything after it
        title = title.split(':', 1)[1].strip()
        return title
        
    # If no colon, check for the "Slide X" pattern
    title_parts = title.split()
    
    # Check if the title starts with "Activity" followed by a number
    if len(title_parts) > 1 and title_parts[0].lower() == "activity" and title_parts[1].replace(":", "").isdigit():
        # Remove the first two parts ("Activity" and the number)
        title_parts = title_parts[2:]
        # Join the remaining parts back into a string
        return " ".join(title_parts).strip()
    
    return title.strip()

def detect_bullet_level(text):
    """
    Detect if text is a bullet point and determine its level.
    Returns a tuple of (is_bullet, level, cleaned_text)
    """
    text = text.strip()
    
    # Check for common bullet point markers
    for marker in BULLET_MARKERS:
        if text.startswith(marker):
            return True, 0, text[len(marker):].strip()
    
    # Check for indented text with bullet markers (sub-bullets)
    if text.startswith('  ') or text.startswith('\\t'):
        # Remove leading whitespace
        stripped = text.lstrip()
        for marker in SUB_BULLET_MARKERS:
            if stripped.startswith(marker):
                return True, 1, stripped[len(marker):].strip()
        
        # If indented but no marker, treat as sub-bullet with no marker
        return True, 1, stripped
    
    # Check for numbered bullets (1., 2., etc.)
    if re.match(r'^\d+\.\s', text):
        return False, 0, text  # Not a bullet but a numbered point
    
    # Not a bullet point
    return False, 0, text

def add_text_box(slide, text, left, top, width, height, font_size=12, bold=False, 
                 italic=False, color=COLORS["text"], alignment=PP_ALIGN.LEFT, 
                 vertical_alignment=MSO_ANCHOR.TOP, level=0, bg_color=None, 
                 border_color=None, shadow=False):
    """Add a text box to a slide with the specified properties."""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = vertical_alignment
    
    # Try to enable auto-size if available
    try:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except:
        print("Exception occurred while enabling auto-size for the text box in the slide")
        pass  # Auto-size not supported in this version
    
    p = text_frame.paragraphs[0]
    p.alignment = alignment
    p.level = level  # Set indentation level
    
    run = p.add_run()
    run.text = text
    
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = color
    
    # Add background color if specified
    if bg_color:
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    
    # Add border if specified
    if border_color:
        line = textbox.line
        line.color.rgb = border_color
        line.width = Pt(1)
    
    # Add shadow if requested
    if shadow and THEME["content_box_shadow"]:
        try:
            # This is a simplified shadow effect
            shadow = textbox.shadow
            shadow.inherit = False
            shadow.visible = True
            shadow.blur_radius = Pt(5)
            shadow.distance = Pt(3)
            shadow.angle = 45
            shadow.color.rgb = RGBColor(0, 0, 0)
            shadow.transparency = 0.7
        except:
            print("Exception occurred while adding shadow to the text box in the slide")
            pass  # Shadow not supported in this version
    
    return textbox

def add_bullet_point(paragraph, text, level=0, font_size=12, bold=False, 
                    italic=False, color=COLORS["text"]):
    """Add a bullet point to an existing paragraph with proper formatting."""
    # Set bullet properties
    paragraph.level = level
    
    # Set bullet visibility
    try:
        paragraph.bullet.visible = True
    except:
        print("Exception occurred while setting bullet visibility in the slide")
        pass  # Bullet customization not supported in this version
    
    # Add the text
    run = paragraph.add_run()
    run.text = text
    
    # Format the text
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = color
    
    return paragraph

def add_bullet_text(slide, text, left, top, width, height, font_size=12, bold=False, 
                   italic=False, color=COLORS["text"], level=0, 
                   bullet_color=None, modern_bullet=False):
    """Add bulleted text to a slide with proper indentation and styling."""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    # Create the paragraph and set its properties
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    
    # Set bullet properties
    p.level = level
    
    # Use modern bullets if enabled
    if THEME["modern_bullets"] and modern_bullet:
        if level > 1:
            text = f"◆ {text}"  # Diamond bullet for sub-sub-points
        elif level > 0:
            text = f"◦ {text}"  # Circle bullet for sub-points
        else:
            text = f"• {text}"  # Bullet for main points
    else:
        # Traditional bullets handled by PowerPoint
        try:
            p.bullet.visible = True
            if bullet_color:
                p.bullet.color.rgb = bullet_color
        except:
            print("Exception occurred while setting bullet properties in the slide")
            pass  # Bullet customization not supported in this version
    
    run = p.add_run()
    run.text = text
    
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = color
    
    return textbox

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, 
              line_color=None, line_width=None, shadow=False, opacity=1.0):
    """Add a shape to a slide with the specified properties."""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        
        # Set transparency if opacity is less than 1
        if opacity < 1.0:
            try:
                shape.fill.transparency = 1.0 - opacity
            except:
                print("Exception occurred while setting transparency for the slide shape")
                pass  # Transparency not supported in this version
    
    if line_color:
        shape.line.color.rgb = line_color
    
    if line_width is not None:
        shape.line.width = Pt(line_width)
    
    # Add shadow if requested
    if shadow:
        try:
            shadow = shape.shadow
            shadow.inherit = False
            shadow.visible = True
            shadow.blur_radius = Pt(5)
            shadow.distance = Pt(3)
            shadow.angle = 45
            shadow.color.rgb = RGBColor(0, 0, 0)
            shadow.transparency = 0.7
        except:
            print("Exception occurred while adding shadow to the slide shape")
            pass  # Shadow not supported in this version
    
    return shape

def add_gradient_background(prs, slide, start_color, end_color, angle=90):
    """Add a gradient background to a slide."""
    # Get slide dimensions from the Presentation object
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Add a rectangle that covers the entire slide
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height
    )
    
    # Remove outline
    shape.line.fill.background()
    
    # Try to set gradient fill
    try:
        fill = shape.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = start_color
        fill.gradient_stops[0].position = 0
        fill.gradient_stops[1].color.rgb = end_color
        fill.gradient_stops[1].position = 1
        fill.gradient_angle = angle
    except:
        # Fallback to solid fill if gradient not supported
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = start_color
    
    # Send to back so it doesn't cover other elements
    try:
        shape.z_order = -100  # Send to back
    except:
        print("Exception occurred while setting z-order for gradient background in the slide")
        pass  # z-order not supported in this version
    
    return shape

def add_corner_accent(slide, color=COLORS["accent1"], size=1.0, position="top-right"):
    """Add a decorative corner accent to a slide."""
    if position == "top-right":
        left = SLIDE_WIDTH - size
        top = 0
    elif position == "top-left":
        left = 0
        top = 0
    elif position == "bottom-right":
        left = SLIDE_WIDTH - size
        top = SLIDE_HEIGHT - size
    elif position == "bottom-left":
        left = 0
        top = SLIDE_HEIGHT - size
    
    # Create a triangle shape for the corner
    points = []
    if position == "top-right":
        points = [(0, 0), (size, 0), (size, size)]
    elif position == "top-left":
        points = [(0, 0), (size, 0), (0, size)]
    elif position == "bottom-right":
        points = [(size, 0), (size, size), (0, size)]
    elif position == "bottom-left":
        points = [(0, 0), (size, size), (0, size)]
    
    # Add a custom shape for the corner accent
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, 
                                  Inches(left), Inches(top), 
                                  Inches(size), Inches(size))
    
    # Set fill color
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    
    # Remove outline
    shape.line.fill.background()
    
    # Set transparency
    try:
        shape.fill.transparency = 0.3  # 70% opacity
    except:
        print("Exception occurred while setting transparency for corner accent in the slide")
        pass  # Transparency not supported in this version
    
    return shape

def add_picture(slide, image_path, left, top, width=None, height=None, 
                shadow=False, border_color=None, border_width=None):
    """Add a picture to a slide with the specified properties."""
    if width and height:
        pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
    elif width:
        pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width))
    elif height:
        pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), height=Inches(height))
    else:
        pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top))
    
    # Add border if specified
    if border_color:
        pic.line.color.rgb = border_color
        if border_width:
            pic.line.width = Pt(border_width)
        else:
            pic.line.width = Pt(1)
    
    # Add shadow if requested
    if shadow:
        try:
            shadow = pic.shadow
            shadow.inherit = False
            shadow.visible = True
            shadow.blur_radius = Pt(5)
            shadow.distance = Pt(3)
            shadow.angle = 45
            shadow.color.rgb = RGBColor(0, 0, 0)
            shadow.transparency = 0.7
        except:
            print("Exception occurred while adding shadow to picture in the slide")
            pass  # Shadow not supported in this version
    
    return pic

def add_table(slide, rows, cols, left, top, width, height, data=None, 
              header_bg_color=None, alt_row_bg_color=None, border_color=None):
    """Add a table to a slide with the specified properties."""
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    
    # Set border color if specified
    if border_color:
        try:
            for row in table.rows:
                for cell in row.cells:
                    cell.border.color.rgb = border_color
        except:
            print("Exception occurred while setting border color in the slide")
            pass  # Border customization not supported in this version
    
    if data:
        for i, row_data in enumerate(data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                
                # Apply header background color to first row
                if i == 0 and header_bg_color:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_bg_color
                # Apply alternating row background color
                elif i > 0 and alt_row_bg_color and i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_row_bg_color
                
                if isinstance(cell_data, dict):
                    text = cell_data.get("text", "")
                    options = cell_data.get("options", {})
                    
                    p = cell.text_frame.paragraphs[0]
                    p.text = text
                    
                    if "fill" in options and "color" in options["fill"]:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = options["fill"]["color"]
                    
                    if "color" in options:
                        p.runs[0].font.color.rgb = options["color"]
                    
                    if "fontSize" in options:
                        p.runs[0].font.size = Pt(options["fontSize"])
                    
                    if "bold" in options:
                        p.runs[0].font.bold = options["bold"]
                    
                    if "italic" in options:
                        p.runs[0].font.italic = options["italic"]
                    
                    if "align" in options:
                        if options["align"] == "center":
                            p.alignment = PP_ALIGN.CENTER
                        elif options["align"] == "right":
                            p.alignment = PP_ALIGN.RIGHT
                else:
                    cell.text = str(cell_data)
    
    return table

def add_footer(slide, title_text, slide_number, total_slides, style="modern"):
    """Add standardized footer to a slide with only the current page number."""
    if style == "modern":
        # Add a subtle divider line above the footer
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.5, FOOTER_Y - 0.05, 9.0, 0.01, 
                 fill_color=COLORS["primary_light"], opacity=0.5)
        
        # Add footer with presentation title
        title_box = add_text_box(slide, title_text, 0.5, FOOTER_Y, 8.5, 0.3, 
                               font_size=10, color=COLORS["primary"], italic=True)
        
        # Add slide number in bottom right - only the current page number
        number_box = add_text_box(slide, f"{slide_number}", 9.0, FOOTER_Y, 0.5, 0.3, 
                                font_size=10, color=COLORS["primary"], 
                                alignment=PP_ALIGN.RIGHT)
    else:
        # Classic footer style
        add_text_box(slide, title_text, 0.5, FOOTER_Y, 8.5, 0.3, font_size=10, color=COLORS["royal_blue"])
        add_text_box(slide, f"{slide_number}", 9.0, FOOTER_Y, 0.5, 0.3, font_size=10, 
                    color=COLORS["text"], alignment=PP_ALIGN.RIGHT)

def calculate_dynamic_spacing(content_items, available_height=AVAILABLE_CONTENT_HEIGHT, min_height=0.4):
    """Calculate dynamic spacing between content items based on available height."""
    if not content_items:
        return min_height
    
    # Calculate spacing based on number of items and available height
    # Ensure minimum spacing and adjust for available space
    item_count = len(content_items)
    spacing = min(0.8, max(min_height, available_height / max(item_count, 1)))
    
    return spacing

def estimate_text_height(text, font_size, width):
    """Estimate the height needed for text based on content length and width."""
    # Approximate characters per line based on font size and width
    # This is a rough estimate - actual rendering may vary
    chars_per_inch = 120 / (font_size / 10)  # Adjust based on average character width
    chars_per_line = int(chars_per_inch * width)
    
    # Calculate number of lines needed
    if chars_per_line <= 0:
        chars_per_line = 1  # Avoid division by zero
    
    text_length = len(text)
    lines = math.ceil(text_length / chars_per_line)
    
    # Calculate height based on lines and font size
    # Add some padding for line spacing
    line_height = (font_size / 72) * 1.2  # Convert points to inches with 1.2 line spacing
    
    # Ensure minimum height
    return max(0.2, lines * line_height)

def check_content_overflow(y_position, content_height, footer_position=FOOTER_Y):
    """Check if content would overflow the slide boundaries."""
    return (y_position + content_height) > (footer_position - 0.2)  # 0.2 inch margin

def create_title_slide(prs, content):
    """Create an enhanced title slide with visual elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background if enabled
    if THEME["use_gradients"]:
        add_gradient_background(prs, slide, COLORS["gradient_start"], COLORS["gradient_end"], angle=135)
    else:
        # Add solid color background
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, 
                 fill_color=COLORS["primary_dark"])
    
    # Add decorative corner accents if enabled
    if THEME["corner_accent"]:
        add_corner_accent(slide, COLORS["accent1"], 2.0, "top-right")
        add_corner_accent(slide, COLORS["accent2"], 1.5, "bottom-left")
    
    # Add title with enhanced styling
    title = content.get("title", "Untitled Presentation")
    title_box = add_text_box(slide, title, 0.5, 1.5, 9, 1.5, font_size=48, 
                           bold=True, color=COLORS["text_light"], 
                           alignment=PP_ALIGN.CENTER, shadow=True)
    
    # Add subtitle with course type and level
    content_type_map = {
        "lecture": "Lecture",
        "tutorial": "Tutorial",
        "workshop": "Workshop",
    }
    
    difficulty_map = {
        "introductory": "Introductory Level",
        "intermediate": "Intermediate Level",
        "advanced": "Advanced Level",
    }
    
    content_type = content.get("contentType", "lecture")
    difficulty_level = content.get("difficultyLevel", "intermediate")
    
    content_type_display = content_type_map.get(content_type, "Lecture")
    difficulty_display = difficulty_map.get(difficulty_level, "Intermediate Level")
    
    subtitle = f"{content_type_display} | {difficulty_display}"
    
    # Add a decorative line between title and subtitle
    line_y = 3.2
    add_shape(slide, MSO_SHAPE.RECTANGLE, 3.5, line_y, 3.0, 0.02, 
             fill_color=COLORS["accent2"])
    
    # Add subtitle text
    subtitle_box = add_text_box(slide, subtitle, 0.5, line_y + 0.2, 9, 0.5, 
                              font_size=28, italic=True, color=COLORS["text_light"], 
                              alignment=PP_ALIGN.CENTER)
    
    return slide

def create_agenda_slide(prs, content, total_slides):
    """Create enhanced agenda slides with visual improvements and overflow handling."""
    # Collect all slide titles and categorize them
    agenda_items = []
    
    # Introduction section
    intro_items = [
        "Learning Outcomes",
        "Key Terms & Concepts"
    ]
    agenda_items.append({"title": "Introduction", "items": intro_items})
    
    # Content slides section
    content_slides = []
    for slide_content in content.get("slides", []):
        title = clean_slide_title(slide_content.get("title", ""))
        if title:
            content_slides.append(title)
    
    if content_slides:
        agenda_items.append({"title": "Main Content", "items": content_slides})
    
    # Activities section
    activities = []
    for idx, activity in enumerate(content.get("activities", [])):
        activity_title = activity.get("title", "")
        activities.append(activity_title)
    
    if activities:
        agenda_items.append({"title": "Activities", "items": activities})
    
    # Test Your Knowledge section (quizzes and discussions)
    knowledge_items = []
    
    # Count quiz and discussion questions
    quiz_count = 0
    discussion_count = 0
    
    for idea in content.get("assessmentIdeas", []):
        idea_type = idea.get("type", "").lower()
        if "quiz" in idea_type:
            quiz_count += len([q for q in idea.get("exampleQuestions", []) if q.get("options")])
        elif "discussion" in idea_type:
            discussion_count += len(idea.get("exampleQuestions", []) if idea.get("exampleQuestions") else [])
    
    if quiz_count > 0:
        knowledge_items.append(f"Quiz Questions ({quiz_count})")
    
    if discussion_count > 0:
        knowledge_items.append(f"Discussion Questions ({discussion_count})")
    
    if knowledge_items:
        agenda_items.append({"title": "Test Your Knowledge", "items": knowledge_items})
    
    # Further Reading section
    if content.get("furtherReadings", []):
        agenda_items.append({"title": "Additional Resources", "items": ["Further Readings & Resources"]})
    
    section_height = 0.5
    item_height = 0.35
    
    total_height_needed = 0
    for section in agenda_items:
        total_height_needed += section_height
        total_height_needed += len(section["items"]) * item_height
    
    available_height = FOOTER_Y - 1.2
    
    slides_needed = math.ceil(total_height_needed / available_height)
    
    # Create agenda slides
    agenda_slides = []
    
    for slide_idx in range(slides_needed):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        agenda_slides.append(slide)
        
        # Add blue header bar
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.8, 
                 fill_color=COLORS["primary"])
        
        # Add title
        title = "Agenda"
        if slide_idx > 0:
            title += f" (continued {slide_idx + 1}/{slides_needed})"
            
        add_text_box(slide, title, 0.5, 0.1, 9, 0.6, 
                   font_size=36, bold=True, color=COLORS["text_light"])
        
        # Create a rounded container for the agenda content
        container = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                            0.3, 0.9, 9.4, FOOTER_Y - 1.1, 
                            fill_color=COLORS["light"], opacity=0.9,
                            line_color=COLORS["primary_light"], line_width=1)
        
        # Distribute content across slides
        y_position = 1.1
        max_y = FOOTER_Y - 0.3  # Leave space for footer
        
        # Track which sections and items we've already displayed
        section_start_idx = 0
        item_start_idx = 0
        
        # For each slide, continue where we left off
        for slide_content_idx in range(slide_idx):
            # Skip through sections until we find where we left off
            for section_idx, section in enumerate(agenda_items[section_start_idx:], section_start_idx):
                section_height = 0.5
                items_height = len(section["items"][item_start_idx:]) * item_height
                total_section_height = section_height + items_height
                
                if total_section_height <= available_height:
                    # Entire section fits, move to next section
                    available_height -= total_section_height
                    section_start_idx += 1
                    item_start_idx = 0
                else:
                    # Section doesn't fit entirely
                    # Calculate how many items fit
                    items_that_fit = int(available_height / item_height)
                    if items_that_fit <= 0:
                        # Not even one item fits, move to next slide
                        available_height = FOOTER_Y - 1.2
                        break
                    
                    # Skip these items for the next slide
                    item_start_idx += items_that_fit
                    available_height = FOOTER_Y - 1.2
                    break
            
            # Reset available height for next slide calculation
            available_height = FOOTER_Y - 1.2
        
        # Now render the content for the current slide
        for section_idx, section in enumerate(agenda_items[section_start_idx:], section_start_idx):
            # Check if we have room for the section title
            if y_position + section_height > max_y:
                break
            
            # Add section title with accent color
            section_title = section["title"]
            add_text_box(slide, section_title, 0.7, y_position, 8.5, section_height, 
                       font_size=24, bold=True, color=COLORS["primary"])
            y_position += section_height
            
            # Add section items as sub-bullets
            for item_idx, item in enumerate(section["items"][item_start_idx:]):
                # Check if we have room for this item
                if y_position + item_height > max_y:
                    # No more room on this slide
                    break
                
                # Create a text box for the sub-bullet
                item_textbox = slide.shapes.add_textbox(
                    Inches(1.0), 
                    Inches(y_position), 
                    Inches(8.0), 
                    Inches(item_height)
                )
                
                # Add the text with proper indentation
                p = item_textbox.text_frame.paragraphs[0]
                p.level = 1  # Set as sub-bullet
                
                # Set bullet visibility
                try:
                    p.bullet.visible = True
                except:
                    # If bullet customization not supported, use text bullet
                    item = f"• {item}"
                
                run = p.add_run()
                run.text = item
                
                # Format the text
                font = run.font
                font.size = Pt(18)
                font.color.rgb = COLORS["text"]
                
                y_position += item_height
            
            # If we've displayed all items in this section, reset item_start_idx for the next section
            if item_idx + item_start_idx >= len(section["items"]) - 1:
                item_start_idx = 0
                section_start_idx += 1
            else:
                # Otherwise, remember where we left off in this section
                item_start_idx += item_idx + 1
                break
        
        # Add footer with slide number
        add_footer(slide, content.get("title", "Untitled Presentation"), 
                 slide_idx + 2, total_slides, THEME["footer_style"])
    
    return agenda_slides

def create_learning_outcomes_slide(prs, content, total_slides):
    """
    Create enhanced learning outcomes slide with visual elements.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add blue header bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.8, fill_color=COLORS["primary"])

    # Add title
    add_text_box(slide, "Learning Outcomes", 0.5, 0.1, 9, 0.6, font_size=36, bold=True, color=COLORS["text_light"])

    # Add intro text
    content_type_map = {
        "lecture": "lecture",
        "tutorial": "tutorial",
        "workshop": "workshop",
    }
    content_type = content.get("contentType", "lecture")
    content_type_display = content_type_map.get(content_type, "lecture")

    intro_text = f"By the end of this {content_type_display}, you will be able to:"
    add_text_box(slide, intro_text, 0.5, 1.0, 9, 0.5, font_size=20, italic=True, color=COLORS["dark"])

    # Add rounded rectangle container for learning outcomes
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.3, 1.7, 9.4, 3.0, fill_color=COLORS["light"],
              opacity=0.9, line_color=COLORS["primary_light"], line_width=1)

    # Add learning outcomes with colored squares and numbering
    learning_outcomes = content.get("learningOutcomes", [])
    y_position = 2.0

    # Define colors for the bullet points
    bullet_colors = [COLORS["emerald"], COLORS["medium_purple"], COLORS["emerald"]]

    for idx, outcome in enumerate(learning_outcomes):
        # Remove existing numbering from the outcome text
        cleaned_outcome = re.sub(r"^\d+\.\s*", "", outcome)

        # Add colored square bullet
        bullet_color = bullet_colors[idx % len(bullet_colors)]
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.7, y_position, 0.15, 0.15, fill_color=bullet_color)

        # Add outcome text slightly higher to align with the bullet center
        add_text_box(slide, cleaned_outcome, 1.0, y_position - 0.125, 8.5, 0.4, font_size=20,
                     color=COLORS["text"], vertical_alignment=MSO_ANCHOR.MIDDLE)

        y_position += 0.6

    # Add footer
    add_footer(slide, content.get("title", "Untitled Presentation"), 1, total_slides, THEME["footer_style"])

    return slide

def create_key_terms_slide(prs, content, total_slides):
    """Create enhanced key terms slides with visual improvements."""
    key_terms = content.get("keyTerms", [])
    if not key_terms:
        return []
    
    # Calculate how many slides we need (4 terms per slide)
    terms_per_slide = 4
    total_terms = len(key_terms)
    slides_needed = (total_terms + terms_per_slide - 1) // terms_per_slide  # Ceiling division
    
    slides = []
    
    for slide_idx in range(slides_needed):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        slides.append(slide)
        
        # Add gradient header if enabled
        add_gradient_background(prs, slide, COLORS["primary"], COLORS["primary_dark"], angle=0)
        
        # Add decorative corner accent if enabled
        add_corner_accent(slide, COLORS["accent3"], 1.0, "bottom-left")
        
        # Add title
        title = "Key Terms & Concepts"
        if slide_idx > 0:
            title += " (continued)"
        add_text_box(slide, title, 0.5, 0.1, 9, 0.6, 
                   font_size=36, bold=True, color=COLORS["text_light"])
        
        # Get terms for this slide
        start_idx = slide_idx * terms_per_slide
        end_idx = min(start_idx + terms_per_slide, total_terms)
        terms_for_slide = key_terms[start_idx:end_idx]
        
        # Calculate the height of the table
        table_height = min(3.5, 0.8 * (len(terms_for_slide) + 1))  # Limit table height to leave space for footer
        
        # Add a background container for the table
        if THEME["content_box_shadow"]:
            add_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, 0.9, 9.2, table_height + 0.2, 
                     fill_color=COLORS["light_alt"], shadow=True)
        
        # Create the table
        table = add_table(slide, len(terms_for_slide) + 1, 2, 0.5, 1.0, 9.0, table_height,
                          header_bg_color=COLORS["royal_blue"],
                          alt_row_bg_color=COLORS["light"],
                          border_color=COLORS["primary_light"])
        
        # Add header row
        header_cells = [
            {"text": "Term", "options": {"fill": {"color": COLORS["royal_blue"]}, "color": COLORS["text_light"], "fontSize": 18, "bold": True, "align": "center"}},
            {"text": "Definition", "options": {"fill": {"color": COLORS["royal_blue"]}, "color": COLORS["text_light"], "fontSize": 18, "bold": True, "align": "center"}}
        ]
        
        for j, cell_data in enumerate(header_cells):
            cell = table.cell(0, j)
            text = cell_data["text"]
            options = cell_data["options"]
            
            p = cell.text_frame.paragraphs[0]
            p.text = text
            
            if "fill" in options and "color" in options["fill"]:
                cell.fill.solid()
                cell.fill.fore_color.rgb = options["fill"]["color"]
            
            if "color" in options:
                p.runs[0].font.color.rgb = options["color"]
            
            if "fontSize" in options:
                p.runs[0].font.size = Pt(options["fontSize"])
            
            if "bold" in options:
                p.runs[0].font.bold = options["bold"]
            
            if "align" in options:
                if options["align"] == "center":
                    p.alignment = PP_ALIGN.CENTER
        
        # Add term rows
        for i, term in enumerate(terms_for_slide):
            row_idx = i + 1  # +1 because of header row
            is_even_row = i % 2 == 0
            row_bg_color = COLORS["background"] if is_even_row else COLORS["light"]
            
            # Term cell
            term_cell = table.cell(row_idx, 0)
            term_cell.text = term.get("term", "")
            term_cell.fill.solid()
            term_cell.fill.fore_color.rgb = row_bg_color
            term_cell.text_frame.paragraphs[0].runs[0].font.bold = True
            term_cell.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
            term_cell.text_frame.paragraphs[0].runs[0].font.color.rgb = COLORS["primary_dark"]
            
            # Definition cell
            def_cell = table.cell(row_idx, 1)
            def_cell.text = term.get("definition", "")
            def_cell.fill.solid()
            def_cell.fill.fore_color.rgb = row_bg_color
            def_cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
        
        # Add footer with white title color
        add_footer(slide, content.get("title", "Untitled Presentation"), slide_idx + 2, total_slides, THEME["footer_style"])
        # Explicitly set the footer title color to white
        add_footer(slide, content.get("title", "Untitled Presentation"), slide_idx + 2, total_slides, THEME["footer_style"])
    
    return slides

def create_content_slides(prs, content, total_slides):
    """Create enhanced slides for the main content with improved formatting and visual elements."""
    slides = content.get("slides", [])
    if not slides:
        return []
    
    result_slides = []
    
    # Calculate starting slide number
    slide_count_offset = 2 + len(content.get("keyTerms", [])) // 4  # Learning Outcomes + Key Terms slides
    
    # Skip the second slide in the content slides section
    for slide_idx, slide_content in enumerate(slides):
        # Skip the second slide (index 1)
        if slide_idx == 1:
            continue
            
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        result_slides.append(slide)
        
        # Add gradient header if enabled
        if THEME["use_gradients"]:
            add_gradient_background(prs, slide, COLORS["primary"], COLORS["primary_dark"], angle=0)
            # Add a semi-transparent white overlay for the content area
            add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0.8, SLIDE_WIDTH, SLIDE_HEIGHT - 0.8, 
                     fill_color=COLORS["background"], opacity=0.9)
        else:
            # Add solid color header bar
            add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.8, 
                     fill_color=COLORS["royal_blue"])
        
        # Add decorative corner accent if enabled
        if THEME["corner_accent"]:
            accent_idx = slide_idx % 3  # Cycle through 3 accent colors
            accent_color = [COLORS["accent1"], COLORS["accent2"], COLORS["accent3"]][accent_idx]
            add_corner_accent(slide, accent_color, 1.0, "bottom-right")
        
        # Add slide title - Remove slide numbers from titles
        original_title = slide_content.get("title", "")
        cleaned_title = clean_slide_title(original_title)
        add_text_box(slide, cleaned_title, 0.5, 0.1, 9, 0.6, 
                   font_size=32, bold=True, color=COLORS["text_light"])
        
        # Process slide content with improved bullet formatting
        content_points = slide_content.get("content", [])
        
        # Create a content container with shadow if enabled
        if THEME["content_box_shadow"]:
            content_height = FOOTER_Y - CONTENT_START_Y - 0.2
            content_container = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                                        0.3, CONTENT_START_Y - 0.1, 9.4, content_height, 
                                        fill_color=COLORS["light_alt"], opacity=0.7,
                                        line_color=COLORS["primary_light"], line_width=1,
                                        shadow=True)
        
        # Create a single text box for all content to ensure proper bullet formatting
        content_textbox = slide.shapes.add_textbox(
            Inches(MAIN_BULLET_INDENT), 
            Inches(CONTENT_START_Y), 
            Inches(9 - MAIN_BULLET_INDENT), 
            Inches(FOOTER_Y - CONTENT_START_Y - 0.3)
        )
        text_frame = content_textbox.text_frame
        text_frame.word_wrap = True
        
        # Check if there are any sub-bullets in the content
        has_sub_bullets = False
        for point in content_points:
            point_text = point if isinstance(point, str) else json.dumps(point)
            if point_text.strip().startswith('  ') or point_text.strip().startswith('\\t') or point_text.strip().startswith('-'):
                has_sub_bullets = True
                break
        
        # Enhanced bullet point detection and formatting
        current_paragraph = text_frame.paragraphs[0]
        is_first_paragraph = True
        
        for point_idx, point in enumerate(content_points):
            point_text = point if isinstance(point, str) else json.dumps(point)
            
            # Determine if this is a bullet point and its level
            is_bullet = False
            bullet_level = 0
            
            # Check for bullet point indicators
            if point_text.strip().startswith('•') or point_text.strip().startswith('*'):
                is_bullet = True
                point_text = point_text.strip()[1:].strip()  # Remove the bullet character
            elif point_text.strip().startswith('-'):
                is_bullet = True
                bullet_level = 1  # Sub-bullet
                point_text = point_text.strip()[1:].strip()  # Remove the hyphen
            elif point_text.strip().startswith('  ') or point_text.strip().startswith('\\t'):
                is_bullet = True
                bullet_level = 1  # Indented text as sub-bullet
                point_text_level = 1  # Indented text as sub-bullet
                point_text = point_text.strip()
            elif not has_sub_bullets:
                # If no sub-bullets exist in the slide, make all points bullets
                is_bullet = True
            
            # Create a new paragraph for each point (except the first one)
            if not is_first_paragraph:
                current_paragraph = text_frame.add_paragraph()
            else:
                is_first_paragraph = False
            
            # Set up bullet formatting
            if is_bullet:
                current_paragraph.level = bullet_level
                try:
                    current_paragraph.bullet.visible = True
                except:
                    # If bullet customization is not supported, use text bullets
                    if THEME["modern_bullets"]:
                        if bullet_level > 0:
                            point_text = f"◦ {point_text}"  # Circle bullet for sub-points
                        else:
                            point_text = f"• {point_text}"  # Bullet for main points
                
            # Add the text
            run = current_paragraph.add_run()
            run.text = point_text
            
            # Format the text based on level
            font = run.font
            if bullet_level == 0:
                font.size = Pt(18)
                font.bold = not is_bullet  # Bold for headings, not for bullets
            else:
                font.size = Pt(16)
                font.bold = False
            
            font.color.rgb = COLORS["text"]
        
        # Add speaker notes to the PowerPoint notes section, not on the slide
        notes = slide_content.get("notes", "")
        if notes:
            notes_text = notes if isinstance(notes, str) else json.dumps(notes)
            
            # Format the notes for better readability
            formatted_notes = notes_text
            
            # Add to PowerPoint's built-in notes section
            if not slide.has_notes_slide:
                slide.notes_slide
            slide.notes_slide.notes_text_frame.text = formatted_notes
        
        # Add footer
        # Adjust slide number to account for skipped slide
        adjusted_slide_idx = slide_idx if slide_idx < 1 else slide_idx - 1
        slide_number = slide_count_offset + adjusted_slide_idx + 1
        add_footer(slide, content.get("title", "Untitled Presentation"), 
                 slide_number, total_slides, THEME["footer_style"])
    
    return result_slides

def extract_facilitation_content(text):
    """
    Extract facilitation notes and learning objectives from text.
    Returns a tuple of (clean_description, facilitation_notes, learning_objectives)
    """
    clean_description = text
    facilitation_notes = ""
    learning_objectives = ""
    
    # Extract facilitation notes with pattern matching
    facilitation_patterns = [
        "Facilitation notes:", 
        "Facilitation Notes:", 
        "FACILITATION NOTES:", 
        "Facilitation notes:", 
        "Facilitation Notes:",
        "Facilitator notes:",
        "Facilitator guidance:",
        "Facilitation tip:"
    ]

    for pattern in facilitation_patterns:
        if pattern in text:
            parts = text.split(pattern, 1)
            clean_description = parts[0].strip()
            facilitation_notes = "Facilitation Notes: " + parts[1].strip()
            break
    
    # Extract learning objectives with pattern matching
    learning_patterns = [
        "Learning Objective:", 
        "Learning Objectives:", 
        "LEARNING OBJECTIVES:", 
        "Learning Objective:", 
        "Learning Objectives:",
        "Success criteria:"
    ]

    for pattern in learning_patterns:
        if pattern in text:
            # If we already extracted facilitation notes, use the clean description
            search_text = clean_description if facilitation_notes else text
            
            if pattern in search_text:
                parts = search_text.split(pattern, 1)
                clean_description = parts[0].strip()
                learning_objectives = "Learning Objective: " + parts[1].strip()
    
    return (clean_description, facilitation_notes, learning_objectives)

def create_activity_slides(prs, content, total_slides):
    """Create activity slides matching the provided design with materials slides for each activity."""
    activities = content.get("activities", [])
    if not activities:
        return []
    
    slides = []
    
    # Calculate starting slide number
    slide_count_offset = (
        2  
        + len(content.get("keyTerms", [])) // 4 - 1  # Additional Key Terms slides
        + len([s for i, s in enumerate(content.get("slides", [])) if i != 1])  # Content slides (excluding second slide)
        + 1
    )
    
    for act_idx, activity in enumerate(activities):
        # Get clean activity title (remove any existing "Activity:" prefix)
        original_title = activity.get("title", "Optimizing Neural Networks")
        clean_title = clean_activity_title(original_title)
        
        # Create the main activity slide
        main_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        slides.append(main_slide)

        # Create the materials slide for EVERY activity (initialize early)
        materials_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        slides.append(materials_slide)
        
        # Add blue header bar
        add_shape(main_slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.8, 
                 fill_color=COLORS["activity_blue"])
        
        # Add activity title with consistent formatting
        title_text = f"Activity {act_idx + 1}: {clean_title}"
        add_text_box(main_slide, title_text, 0.5, 0.07, 9.0, 0.6, 
                   font_size=22, bold=True, color=COLORS["text_light"])
        
        # Add purple badge for type and duration
        badge_shape = add_shape(main_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                              0.5, 0.9, 5.0, 0.5, 
                              fill_color=COLORS["activity_purple"], 
                              line_color=None)
        
        # Add type and duration text
        activity_type = activity.get("type", "Exercise")
        activity_duration = activity.get("duration", "20 minutes")
        type_duration_text = f"Type: {activity_type} | Duration: {activity_duration}"
        add_text_box(main_slide, type_duration_text, 0.6, 0.9, 4.8, 0.5, 
                   font_size=16, italic=True, color=COLORS["text_light"], 
                   vertical_alignment=MSO_ANCHOR.MIDDLE)
        
        # Create a large rounded container for all content
        content_container = add_shape(main_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    0.5, 1.5, 9.0, 3.5, 
                                    fill_color=COLORS["background"], 
                                    line_color=COLORS["activity_purple"],
                                    line_width=1)
        
        # Add activity description
        activity_description = activity.get("description", "Optimize a neural network model for Gaudi-3 using quantization, pruning, or knowledge distillation to improve performance.")
        
        # Extract facilitation notes and learning objectives
        clean_description, facilitation_notes, learning_objectives = extract_facilitation_content(activity_description)

        # Update description if facilitation notes or learning objectives were found
        if facilitation_notes or learning_objectives:
            # Update the description text boxes with clean description
            add_text_box(main_slide, clean_description, 0.7, 1.7, 8.6, 0.6, 
                       font_size=20, color=COLORS["text"])
            add_text_box(materials_slide, clean_description, 0.7, 1.7, 8.6, 0.6, 
                       font_size=20, color=COLORS["text"])
            
            # Combine notes
            combined_notes = ""
            if facilitation_notes:
                combined_notes += facilitation_notes
            if learning_objectives:
                combined_notes += " " + learning_objectives if combined_notes else learning_objectives
            
            if combined_notes:
                # Add notes to the PowerPoint notes section
                if not main_slide.has_notes_slide:
                    main_slide.notes_slide
                main_slide.notes_slide.notes_text_frame.text = combined_notes
                
                # Also add to materials slide notes
                if not materials_slide.has_notes_slide:
                    materials_slide.notes_slide
                materials_slide.notes_slide.notes_text_frame.text = combined_notes
            
        
        # Use the clean description for the rest of the function
        activity_description = clean_description
        
        # Add blue header bar
        add_shape(materials_slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.8, 
                 fill_color=COLORS["activity_blue"])
        
        # Add activity title with (Materials) suffix
        materials_title = f"Activity {act_idx + 1}: {clean_title} (Materials)"
        add_text_box(materials_slide, materials_title, 0.5, 0.07, 9.0, 0.6, 
                   font_size=22, bold=True, color=COLORS["text_light"])
        
        # Add purple badge for type and duration (same as main slide)
        badge_shape = add_shape(materials_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                              0.5, 0.9, 5.0, 0.5, 
                              fill_color=COLORS["activity_purple"], 
                              line_color=None)
        
        # Add type and duration text
        add_text_box(materials_slide, type_duration_text, 0.6, 0.9, 4.8, 0.5, 
                   font_size=16, italic=True, color=COLORS["text_light"], 
                   vertical_alignment=MSO_ANCHOR.MIDDLE)
        
        # Create a large rounded container for all content
        content_container = add_shape(materials_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    0.5, 1.5, 9.0, 3.5, 
                                    fill_color=COLORS["background"], 
                                    line_color=COLORS["activity_purple"],
                                    line_width=1)
        
        # Add activity description (same as main slide)
        add_text_box(main_slide, activity_description, 0.7, 1.7, 8.6, 0.6, 
                   font_size=20, color=COLORS["text"])
        add_text_box(materials_slide, activity_description, 0.7, 1.7, 8.6, 0.6, 
                   font_size=20, color=COLORS["text"])
        
        # Extract facilitation notes with more robust pattern matching
        facilitation_notes = ""
        description_clean = activity_description
        
        # Check for various facilitation note formats
        facilitation_patterns = [
            "Facilitation notes:", 
            "Facilitation Notes:", 
            "FACILITATION NOTES:", 
            "Facilitation notes:", 
            "Facilitation Notes:",
            "Facilitator notes:",
            "Facilitator guidance:",
            "Success criteria:",
            "Success notes:"
            "Learning Objective:"
        ]
        
        for pattern in facilitation_patterns:
            if pattern in activity_description:
                parts = activity_description.split(pattern, 1)
                description_clean = parts[0].strip()
                facilitation_notes = "Facilitation Notes: " + parts[1].strip()
                break
        
        # If facilitation notes were found, update the slides
        if facilitation_notes:
            # Update the description text boxes with clean description
            add_text_box(main_slide, description_clean, 0.7, 1.7, 8.6, 0.6, 
                       font_size=20, color=COLORS["text"])
            add_text_box(materials_slide, description_clean, 0.7, 1.7, 8.6, 0.6, 
                       font_size=20, color=COLORS["text"])
            
            # Add facilitation notes to the PowerPoint notes section
            if not main_slide.has_notes_slide:
                main_slide.notes_slide
            main_slide.notes_slide.notes_text_frame.text = facilitation_notes
            
            # Also add to materials slide notes
            if not materials_slide.has_notes_slide:
                materials_slide.notes_slide
            materials_slide.notes_slide.notes_text_frame.text = facilitation_notes
            
            # Add a visual indicator that facilitation notes are available
            add_shape(main_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                     8.5, 0.9, 1.0, 0.5, 
                     fill_color=COLORS["activity_green"], 
                     line_color=None)
            add_text_box(main_slide, "Notes Available", 8.6, 0.95, 0.8, 0.4, 
                       font_size=12, bold=True, color=COLORS["text_light"], 
                       alignment=PP_ALIGN.CENTER)
            
            # Add the same indicator to materials slide
            add_shape(materials_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                     8.5, 0.9, 1.0, 0.5, 
                     fill_color=COLORS["activity_green"], 
                     line_color=None)
            add_text_box(materials_slide, "Notes Available", 8.6, 0.95, 0.8, 0.4, 
                       font_size=12, bold=True, color=COLORS["text_light"], 
                       alignment=PP_ALIGN.CENTER)
        else:
            # If no facilitation notes were found, use the original description
            activity_description = description_clean
        
        # Extract learning objectives for tutorial slides if present
        learning_objectives = ""
        isTutorial = activity.get("type", "").lower() == "tutorial"
        if isTutorial and "Learning Objective:" in activity_description:
            parts = activity_description.split("Learning Objective:", 1)
            activity_description = parts[0].strip()
            learning_objectives = "Learning Objective: " + parts[1].strip()
            
            # Update the description text boxes with clean description
            add_text_box(main_slide, activity_description, 0.7, 1.7, 8.6, 0.6, 
                       font_size=20, color=COLORS["text"])
            add_text_box(materials_slide, activity_description, 0.7, 1.7, 8.6, 0.6, 
                       font_size=20, color=COLORS["text"])
            
            # Append learning objectives to notes
            if not main_slide.has_notes_slide:
                main_slide.notes_slide
            current_notes = main_slide.notes_slide.notes_text_frame.text
            main_slide.notes_slide.notes_text_frame.text = current_notes + "  " + learning_objectives if current_notes else learning_objectives
            
            # Also add to materials slide notes
            if not materials_slide.has_notes_slide:
                materials_slide.notes_slide
            current_notes = materials_slide.notes_slide.notes_text_frame.text
            materials_slide.notes_slide.notes_text_frame.text = current_notes + "  " + learning_objectives if current_notes else learning_objectives
        
        # Add instructions section with blue vertical accent bar
        instructions_y = 2.4
        add_shape(main_slide, MSO_SHAPE.RECTANGLE, 0.7, instructions_y + 0.1, 0.1, 2.0, 
                 fill_color=COLORS["activity_blue"])
        
        # Add "Instructions:" heading
        add_text_box(main_slide, "Instructions:", 0.9, instructions_y, 8.3, 0.4, 
                   font_size=22, bold=True, color=COLORS["text"])
        
        # Add numbered instructions
        instructions = activity.get("instructions", [])
        
        # Create a text box for instructions
        instructions_textbox = main_slide.shapes.add_textbox(
            Inches(0.9), 
            Inches(instructions_y + 0.5), 
            Inches(8.3), 
            Inches(1.5)
        )
        instructions_frame = instructions_textbox.text_frame
        instructions_frame.word_wrap = True
        
        # Add each instruction as a numbered point
        for idx, instruction in enumerate(instructions):
            if idx > 0:
                p = instructions_frame.add_paragraph()
            else:
                p = instructions_frame.paragraphs[0]
            
            # Format as numbered list
            instruction_text = instruction if isinstance(instruction, str) else json.dumps(instruction)
            p.text = f"{idx + 1}. {instruction_text}"
            
            # Format the text
            for run in p.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = COLORS["text"]
        
        # Add orange triangle in bottom left corner
        triangle = add_shape(main_slide, MSO_SHAPE.RIGHT_TRIANGLE, 
                           0, SLIDE_HEIGHT - 1.5, 1.5, 1.5, 
                           fill_color=COLORS["activity_orange"])
        
        # Add footer with blue divider line
        add_shape(main_slide, MSO_SHAPE.RECTANGLE, 0.5, FOOTER_Y - 0.05, 9.0, 0.01, 
                 fill_color=COLORS["primary_light"])
        
        # Add presentation title on left side of footer
        presentation_title = content.get("title", "Optimizing Neural Networks on Gaudi-3 AI Accelerator")
        add_text_box(main_slide, presentation_title, 0.5, FOOTER_Y, 8.0, 0.3, 
                   font_size=10, color=COLORS["primary"], italic=True)
        
        # Add slide number on right side
        slide_number = slide_count_offset + (act_idx * 2) + 1  # Each activity has 2 slides
        add_text_box(main_slide, f"{slide_number}", 9.0, FOOTER_Y, 0.5, 0.3, 
                   font_size=10, color=COLORS["primary"], 
                   alignment=PP_ALIGN.RIGHT)
        
        # Add materials section with green vertical accent bar
        materials_y = 2.4
        add_shape(materials_slide, MSO_SHAPE.RECTANGLE, 0.7, materials_y + 0.1, 0.1, 2.0, 
                 fill_color=COLORS["activity_green"])
        
        # Add "Materials needed:" heading
        add_text_box(materials_slide, "Materials needed:", 0.9, materials_y, 8.3, 0.4, 
                   font_size=22, bold=True, color=COLORS["text"])
        
        # Add materials list - dynamically from the backend data
        materials = activity.get("materials", ["Gaudi-3 optimization tools", "Neural network models"])
        
        # Create a text box for materials
        materials_textbox = materials_slide.shapes.add_textbox(
            Inches(0.9), 
            Inches(materials_y + 0.5), 
            Inches(8.3), 
            Inches(1.5)
        )
        materials_frame = materials_textbox.text_frame
        materials_frame.word_wrap = True
        
        # Add each material as a bullet point
        for idx, material in enumerate(materials):
            if idx > 0:
                p = materials_frame.add_paragraph()
            else:
                p = materials_frame.paragraphs[0]
            
            # Format as bullet list
            material_text = material if isinstance(material, str) else json.dumps(material)
            p.text = f"• {material_text}"
            
            # Format the text
            for run in p.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = COLORS["text"]
        
        # Add orange triangle in bottom left corner
        triangle = add_shape(materials_slide, MSO_SHAPE.RIGHT_TRIANGLE, 
                           0, SLIDE_HEIGHT - 1.5, 1.5, 1.5, 
                           fill_color=COLORS["activity_orange"])
        
        # Add footer with blue divider line
        add_shape(materials_slide, MSO_SHAPE.RECTANGLE, 0.5, FOOTER_Y - 0.05, 9.0, 0.01, 
                 fill_color=COLORS["primary_light"])
        
        # Add presentation title on left side of footer
        add_text_box(materials_slide, presentation_title, 0.5, FOOTER_Y, 8.0, 0.3, 
                   font_size=10, color=COLORS["primary"], italic=True)
        
        # Add slide number on right side (increment by 1 from the main slide)
        slide_number = slide_count_offset + (act_idx * 2) + 2  # +2 for the materials slide
        add_text_box(materials_slide, f"{slide_number}", 9.0, FOOTER_Y, 0.5, 0.3, 
                   font_size=10, color=COLORS["primary"], 
                   alignment=PP_ALIGN.RIGHT)
    
    return slides

def create_quiz_slides(prs, content, total_slides):
    """Create enhanced quiz question and answer slides with visual improvements based on the new design."""
    assessment_ideas = content.get("assessmentIdeas", [])
    slides = []
    
    # Calculate slide count offset
    slide_count_offset = (
        2 +  # Learning Outcomes + first Key Terms slide
        len(content.get("keyTerms", [])) // 4 - 1 +  # Additional Key Terms slides
        len([s for i, s in enumerate(content.get("slides", [])) if i != 1]) +  # Content slides (excluding second slide)
        len(content.get("activities", [])) * 2  # Activity slides (2 per activity)
        + 1
    )
    
    quiz_slide_count = 0
    
    for idea_idx, idea in enumerate(assessment_ideas):
        idea_type = idea.get("type", "Assessment")
        is_quiz = "quiz" in idea_type.lower()
        
        if not is_quiz:
            continue
        
        example_questions = idea.get("exampleQuestions", [])
        
        for q_idx, question in enumerate(example_questions):
            question_text = question.get("question", "Example question")
            options = question.get("options", [])
            
            if not options:
                continue
            
            # Create question slide with enhanced styling
            q_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            slides.append(q_slide)
            quiz_slide_count += 1
            
            # Add full-width blue header bar
            add_shape(q_slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 1.0, 
                    fill_color=COLORS["primary"])
            
            # Add slide title in the blue header
            add_text_box(q_slide, f"Quiz Question {q_idx + 1}", 0.5, 0.2, 9.0, 0.6, 
                      font_size=36, bold=True, color=COLORS["text_light"], 
                      alignment=PP_ALIGN.CENTER, vertical_alignment=MSO_ANCHOR.MIDDLE)
            
            # Add the question text with enhanced styling - light gray rounded rectangle
            question_box = add_shape(q_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                                  0.5, 1.1, 9.0, 0.8, 
                                  fill_color=COLORS["light"], 
                                  line_color=COLORS["light"],  # Same color as fill for seamless look
                                  line_width=1, shadow=False)

            # Add the question text
            add_text_box(q_slide, question_text, 0.7, 1.2, 8.6, 0.6, 
                      font_size=20, bold=True, color=COLORS["text"])
            
            # Add options in a 2x2 grid with enhanced styling
            options_per_row = 2
            option_width = 4.3
            option_height = 1.0  # Reduced height for cleaner look
            option_gap = 0.4
            start_y = 2.2  # Adjusted starting position
            
            for opt_idx, option in enumerate(options):
                row = opt_idx // options_per_row
                col = opt_idx % options_per_row
                option_x = 0.5 + col * (option_width + option_gap)
                option_y = start_y + row * (option_height + 0.4)
                
                # Add option box with enhanced styling - light gray rounded rectangle
                option_box = add_shape(q_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    option_x, option_y, option_width, option_height, 
                                    fill_color=COLORS["light"], 
                                    line_color=COLORS["light"],  # Same color as fill for seamless look
                                    line_width=1, shadow=False)
                
                # Add option letter in a blue circle with enhanced styling
                circle_size = 0.6  # Slightly smaller circle
                circle_x = option_x + 0.2
                circle_y = option_y + (option_height - circle_size) / 2  # Center vertically
                
                circle = add_shape(q_slide, MSO_SHAPE.OVAL, 
                                circle_x, circle_y, circle_size, circle_size, 
                                fill_color=COLORS["primary"])  # Blue circle
                
                # Add letter - ensure it's centered in the circle
                letter = chr(65 + opt_idx)  # A, B, C, D...
                letter_textbox = add_text_box(q_slide, letter, circle_x, circle_y, circle_size, circle_size, 
                                          font_size=24, bold=True, color=COLORS["text_light"], 
                                          alignment=PP_ALIGN.CENTER, vertical_alignment=MSO_ANCHOR.MIDDLE)
                
                # Add option text with improved styling - centered vertically and horizontally
                text_x = circle_x + circle_size + 0.2
                text_width = option_width - (text_x - option_x) - 0.2
                
                # Create text box for option text with center alignment
                text_box = add_text_box(q_slide, option, text_x, option_y, 
                                     text_width, option_height, 
                                     font_size=18, color=COLORS["text"],
                                     alignment=PP_ALIGN.CENTER,
                                     vertical_alignment=MSO_ANCHOR.MIDDLE)
            
            # Add footer with presentation title and slide number
            presentation_title = content.get("title", "Untitled Presentation")
            slide_number = slide_count_offset + quiz_slide_count
            
            # Add footer
            add_footer(q_slide, presentation_title, slide_number, total_slides, THEME["footer_style"])
            
            # Create answer slide with enhanced styling to match the image
            a_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            slides.append(a_slide)
            quiz_slide_count += 1
            
            # Add full-width blue header bar (taller than before)
            add_shape(a_slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 1.2, 
                    fill_color=COLORS["primary"])
            
            # Add slide title in the blue header - centered
            add_text_box(a_slide, f"Quiz Answer {q_idx + 1}", 0.5, 0.3, 9.0, 0.6, 
                      font_size=40, bold=True, color=COLORS["text_light"], 
                      alignment=PP_ALIGN.CENTER, vertical_alignment=MSO_ANCHOR.MIDDLE)
            
            # Add the question text as a reminder with enhanced styling - light gray rounded rectangle
            question_box = add_shape(a_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                                  0.5, 1.4, 9.0, 0.8, 
                                  fill_color=COLORS["light"], 
                                  line_color=COLORS["light"],
                                  line_width=1, shadow=True)
            
            add_text_box(a_slide, f"Question: {question_text}", 0.7, 1.5, 8.6, 0.6, 
                      font_size=18, italic=True, color=COLORS["text"])
            
            # Add correct answer with enhanced styling - dark background with green border
            correct_answer = question.get("correctAnswer", "")
            if correct_answer:
                # Add a highlight box for the answer with enhanced styling
                answer_box = add_shape(a_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    0.5, 2.4, 9.0, 1.0, 
                                    fill_color=COLORS["dark_alt"], 
                                    line_color=COLORS["success"], 
                                    line_width=3, shadow=True)
                
                # Add "Correct Answer:" text with enhanced styling - yellow/gold text
                add_text_box(a_slide, "Correct Answer:", 0.7, 2.5, 8.6, 0.4, 
                          font_size=20, bold=True, color=COLORS["warning"])
                
                # Add the answer text with enhanced styling - white text
                add_text_box(a_slide, correct_answer, 0.7, 2.9, 8.6, 0.4, 
                          font_size=18, color=COLORS["text_light"])
            
            # Add explanation section with enhanced styling
            explanation = question.get("explanation", "")
            if explanation:
                # Add explanation container with enhanced styling - light background with subtle border
                explanation_box = add_shape(a_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                                         0.5, 3.6, 9.0, 1.4, 
                                         fill_color=COLORS["light"], 
                                         line_color=COLORS["primary_light"], 
                                         line_width=1, shadow=True)
                
                # Add blue vertical accent bar on the left
                add_shape(a_slide, MSO_SHAPE.RECTANGLE, 0.7, 3.7, 0.1, 1.2, 
                        fill_color=COLORS["primary"])
                
                # Add explanation title
                add_text_box(a_slide, "Explanation:", 0.9, 3.7, 8.5, 0.4, 
                          font_size=20, bold=True, color=COLORS["text"])
                
                # Format and add the explanation with improved styling
                explanation_text = explanation if isinstance(explanation, str) else json.dumps(explanation)
                add_text_box(a_slide, explanation_text, 0.9, 4.2, 8.5, 0.7, 
                          font_size=16, color=COLORS["text"])
            
            # Add footer with presentation title on left and slide number on right
            presentation_title = content.get("title", "Untitled Presentation")
            slide_number = slide_count_offset + quiz_slide_count
            
            # Add divider line above footer
            add_shape(a_slide, MSO_SHAPE.RECTANGLE, 0.5, FOOTER_Y - 0.05, 9.0, 0.01, 
                    fill_color=COLORS["primary_light"])
            
            # Add presentation title on left side of footer
            add_text_box(a_slide, presentation_title, 0.5, FOOTER_Y, 8.0, 0.3, 
                      font_size=10, color=COLORS["primary"], italic=True)
            
            # Add slide number on right side
            add_text_box(a_slide, f"{slide_number}", 9.0, FOOTER_Y, 0.5, 0.3, 
                      font_size=10, color=COLORS["primary"], 
                      alignment=PP_ALIGN.RIGHT)
    
    return slides

def create_discussion_slides(prs, content, total_slides):
    """Create discussion question slides with guidance on the same slide and answer slides."""
    assessment_ideas = content.get("assessmentIdeas", [])
    slides = []
    
    # Calculate slide count offset including quiz slides
    slide_count_offset = (
        2 +  # Learning Outcomes + first Key Terms slide
        len(content.get("keyTerms", [])) // 4 - 1 +  # Additional Key Terms slides
        len([s for i, s in enumerate(content.get("slides", [])) if i != 1]) +  # Content slides (excluding second slide)
        len(content.get("activities", [])) * 2  # Activity slides (2 per activity)
        + 1
    )
    
    # Count quiz slides
    quiz_count = 0
    for idea in content.get("assessmentIdeas", []):
        if "quiz" in idea.get("type", "").lower():
            quiz_count += len([q for q in idea.get("exampleQuestions", []) if q.get("options")]) * 2
    
    discussion_slide_count = 0
    
    for idea_idx, idea in enumerate(assessment_ideas):
        idea_type = idea.get("type", "Assessment")
        is_discussion = "discussion" in idea_type.lower()
        
        if not is_discussion:
            continue
        
        example_questions = idea.get("exampleQuestions", [])
        
        for q_idx, question in enumerate(example_questions):
            question_text = question.get("question", "Example question")
            guidance = question.get("correctAnswer", "")
            
            # Create question slide with guidance
            q_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            slides.append(q_slide)
            discussion_slide_count += 1
            
            # Add header bar
            add_shape(q_slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.8, 
                     fill_color=COLORS["primary"])
            
            # Add slide title
            add_text_box(q_slide, f"Discussion Question {q_idx + 1}", 0.5, 0.1, 9.0, 0.6, 
                       font_size=32, bold=True, color=COLORS["text_light"])
            
            # Add the question text in a rounded box at the top
            question_box = add_shape(q_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                                  0.5, 1.1, 9.0, 0.8, 
                                  fill_color=COLORS["light"], 
                                  line_color=COLORS["light"],
                                  line_width=1, shadow=True)
            
            question_text_box = add_text_box(q_slide, question_text, 0.7, 1.2, 8.6, 0.6, 
                       font_size=20, bold=True, color=COLORS["text"])
            
            # Dynamically calculate the height of the question text
            question_text_height = estimate_text_height(question_text, 20, 8.6)
            next_text_y = 1.2 + question_text_height + 1.2  # Add some padding
            
            # Add discussion prompt dynamically below the question
            add_shape(q_slide, MSO_SHAPE.RECTANGLE, 0.7, next_text_y, 0.1, 0.4, 
                     fill_color=COLORS["primary"])
            add_text_box(q_slide, "Group Discussion:", 0.9, next_text_y, 8.5, 0.4, 
                       font_size=20, bold=True, color=COLORS["primary"])
            
            # Add discussion instructions dynamically below the prompt
            instructions_y = next_text_y + 0.5
            instructions = "Discuss this question with your group and prepare to share your thoughts with the class."
            add_text_box(q_slide, instructions, 0.9, instructions_y, 8.5, 0.4, 
                       font_size=18, color=COLORS["text"])
            
            # Add footer with presentation title and slide number
            presentation_title = content.get("title", "Untitled Presentation")
            slide_number = slide_count_offset + quiz_count + discussion_slide_count
            
            # Add footer
            add_footer(q_slide, presentation_title, slide_number, total_slides, THEME["footer_style"])
            
            # Create answer slide with facilitator guidance
            a_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            slides.append(a_slide)
            discussion_slide_count += 1
            
            # Add full-width blue header bar
            add_shape(a_slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.8, 
                    fill_color=COLORS["primary"])
            
            # Add slide title
            add_text_box(a_slide, f"Facilitator Guidance: Question {q_idx + 1}", 0.5, 0.1, 9.0, 0.6, 
                      font_size=32, bold=True, color=COLORS["text_light"])
            
            # Add the question text as a reminder
            question_box = add_shape(a_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                                  0.5, 1.1, 9.0, 0.8, 
                                  fill_color=COLORS["light"], 
                                  line_color=COLORS["light"],
                                  line_width=1, shadow=True)
            
            add_text_box(a_slide, f"Question: {question_text}", 0.7, 1.2, 8.6, 0.6, 
                      font_size=18, italic=True, color=COLORS["text"])
            
            # Add facilitator guidance section dynamically
            guidance_y = 1.2 + question_text_height + 0.7
            if guidance:
                # Add a guidance container with enhanced styling
                guidance_box = add_shape(a_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                                      0.5, guidance_y, 9.0, 2.8, 
                                      fill_color=COLORS["light"], 
                                      line_color=COLORS["accent2"], 
                                      line_width=2, shadow=True)
                
                # Add accent bar on the left
                add_shape(a_slide, MSO_SHAPE.RECTANGLE, 0.7, guidance_y + 0.1, 0.1, 2.5, 
                        fill_color=COLORS["accent2"])
                
                # Add "Facilitator Guidance:" heading
                add_text_box(a_slide, "Facilitator Guidance:", 0.9, guidance_y + 0.1, 8.5, 0.4, 
                          font_size=20, bold=True, color=COLORS["accent2"])
                
                # Format and add the guidance with improved styling
                guidance_text = guidance if isinstance(guidance, str) else json.dumps(guidance)
                add_text_box(a_slide, guidance_text, 0.9, guidance_y + 0.6, 8.3, 1.5, 
                          font_size=16, color=COLORS["text"])
            
            # Add footer
            slide_number = slide_count_offset + quiz_count + discussion_slide_count
            add_footer(a_slide, presentation_title, slide_number, total_slides, THEME["footer_style"])
    
    return slides

def create_further_readings_slides(prs, content, total_slides):
    """Create further readings slides with version 7 styling."""
    readings = content.get("furtherReadings", [])
    if not readings:
        return []
    
    slides = []
    
    # Calculate slide count offset
    slide_count_offset = (
        2 +  # Learning Outcomes + first Key Terms slide
        len(content.get("keyTerms", [])) // 4 - 1 +  # Additional Key Terms slides
        len([s for i, s in enumerate(content.get("slides", [])) if i != 1]) +  # Content slides (excluding second slide)
        len(content.get("activities", [])) * 2  # Activity slides (2 per activity)
        + 1
    )
    
    # Count quiz slides
    quiz_count = 0
    for idea in content.get("assessmentIdeas", []):
        if "quiz" in idea.get("type", "").lower():
            quiz_count += len([q for q in idea.get("exampleQuestions", []) if q.get("options")]) * 2
    
    # Count discussion slides - now 2 slides per question (with guidance)
    discussion_count = 0
    for idea in content.get("assessmentIdeas", []):
        if "discussion" in idea.get("type", "").lower():
            discussion_count += len(idea.get("exampleQuestions", []) if idea.get("exampleQuestions") else []) * 2
    
    # Calculate how many slides we need (2 readings per slide)
    readings_per_slide = 2
    total_readings = len(readings)
    slides_needed = (total_readings + readings_per_slide - 1) // readings_per_slide
    
    for slide_idx in range(slides_needed):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        slides.append(slide)
        
        # Add header bar (version 7 style - blue header)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.8, 
                 fill_color=COLORS["primary"])
        
        # Add title
        title = "Further Readings & Resources"
        if slide_idx > 0:
            title += f" (continued)"
        add_text_box(slide, title, 0.5, 0.1, 9, 0.6, 
                   font_size=32, bold=True, color=COLORS["text_light"])
        
        # Add content container
        if THEME["content_box_shadow"]:
            content_container = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                                       0.3, 1.0, 9.4, FOOTER_Y - 1.2, 
                                       fill_color=COLORS["light_alt"], opacity=0.7,
                                       line_color=COLORS["primary_light"], line_width=1,
                                       shadow=True)
        
        # Get readings for this slide
        start_idx = slide_idx * readings_per_slide
        end_idx = min(start_idx + readings_per_slide, total_readings)
        readings_for_slide = readings[start_idx:end_idx]
        
        # Calculate the starting Y position for the readings
        y_position = 1.2
        
        # Add readings with version 7 styling
        for i, reading in enumerate(readings_for_slide):
            # Add reading title with accent bar
            reading_title = reading.get("title", "Untitled Reading")
            add_shape(slide, MSO_SHAPE.RECTANGLE, 0.7, y_position, 0.1, 0.4,
                     fill_color=COLORS["primary"])
            add_text_box(slide, reading_title, 0.9, y_position, 8.3, 0.4,
                       font_size=20, bold=True, color=COLORS["primary"])
            y_position += 0.5
            
            # Add reading URL or author
            reading_author = reading.get("author", "")
            reading_description = reading.get("readingDescription", "")
            add_text_box(slide, f"Author: {reading_author}", 0.9, y_position, 8.3, 0.3,
                       font_size=16, italic=True, color=COLORS["primary"])
            y_position += 0.4
            
            # Add reading description
            add_text_box(slide, reading_description, 0.9, y_position, 8.3, 0.6,
                       font_size=16, color=COLORS["text"])
            
            # Add separator line if not the last reading
            if i < len(readings_for_slide) - 1:
                y_position += 0.8
                add_shape(slide, MSO_SHAPE.RECTANGLE, 0.7, y_position, 8.5, 0.01,
                         fill_color=COLORS["primary_light"], opacity=0.5)
                y_position += 0.2
            else:
                y_position += 0.8
        
        # Add footer
        slide_number = slide_count_offset + quiz_count + discussion_count + slide_idx + 1
        add_footer(slide, content.get("title", "Untitled Presentation"), slide_number, total_slides, THEME["footer_style"])
    
    return slides

def create_facilitation_notes_slide(prs, content, total_slides):
    """Create a dedicated slide that summarizes all facilitation notes."""
    activities = content.get("activities", [])

    # Calculate slide count offset
    slide_count_offset = (
        2 +  # Learning Outcomes + first Key Terms slide
        len(content.get("keyTerms", [])) // 4 - 1 +  # Additional Key Terms slides
        len([s for i, s in enumerate(content.get("slides", [])) if i != 1]) +  # Content slides (excluding second slide)
        len(content.get("activities", [])) * 2  # Activity slides (2 per activity)
        + 1
        + len(content.get("furtherReadings", [])) // 2  # Further readings slides
    )

    # Count quiz slides
    quiz_count = 0
    for idea in content.get("assessmentIdeas", []):
        if "quiz" in idea.get("type", "").lower():
            quiz_count += len([q for q in idea.get("exampleQuestions", []) if q.get("options")]) * 2
    
    # Count discussion slides - now 2 slides per question (with guidance)
    discussion_count = 0
    for idea in content.get("assessmentIdeas", []):
        if "discussion" in idea.get("type", "").lower():
            discussion_count += len(idea.get("exampleQuestions", []) if idea.get("exampleQuestions") else []) * 2
    
    # Check if there are any activities with facilitation notes
    has_facilitation_notes = False
    for activity in activities:
        description = activity.get("description", "")
        _, facilitation_notes, _ = extract_facilitation_content(description)
        if facilitation_notes:
            has_facilitation_notes = True
            break
    
    # If no facilitation notes, don't create the slide
    if not has_facilitation_notes:
        return None
    
    # Create the slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add header bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.8, 
             fill_color=COLORS["primary"])
    
    # Add title
    add_text_box(slide, "Facilitation Notes Summary", 0.5, 0.1, 9.0, 0.6, 
               font_size=32, bold=True, color=COLORS["text_light"])
    
    # Create a container for the notes
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
             0.5, 1.0, 9.0, FOOTER_Y - 1.2, 
             fill_color=COLORS["light_alt"], 
             line_color=COLORS["primary_light"],
             line_width=1, shadow=True)
    
    # Add facilitation notes from each activity
    y_position = 1.2
    for idx, activity in enumerate(activities):
        title = activity.get("title", "")
        description = activity.get("description", "")
        
        _, facilitation_notes, _ = extract_facilitation_content(description)
        
        if facilitation_notes:
            # Add activity title
            add_shape(slide, MSO_SHAPE.RECTANGLE, 0.7, y_position, 0.1, 0.4, 
                     fill_color=COLORS["activity_green"])
            add_text_box(slide, title, 0.9, y_position, 8.3, 0.4, 
                       font_size=18, bold=True, color=COLORS["primary"])
            y_position += 0.5
            
            # Add facilitation notes
            notes_text = facilitation_notes.replace("Facilitation Notes: ", "")
            add_text_box(slide, notes_text, 0.9, y_position, 8.3, 0.6, 
                       font_size=14, color=COLORS["text"])
            y_position += 0.8
            
            # Add separator if not the last item
            if idx < len(activities) - 1:
                add_shape(slide, MSO_SHAPE.RECTANGLE, 0.7, y_position, 8.5, 0.01, 
                         fill_color=COLORS["primary_light"], opacity=0.5)
                y_position += 0.3
            
            # Check if we need a new slide (if y_position is too large)
            if y_position > FOOTER_Y - 0.5:
                # Add "continued on next slide" text
                add_text_box(slide, "Continued on next slide...", 0.9, y_position - 0.3, 8.3, 0.3, 
                           font_size=12, italic=True, color=COLORS["text_muted"])
                
                # Add footer
                add_footer(slide, content.get("title", "Untitled Presentation"), 
                         total_slides - 1, total_slides + 1, THEME["footer_style"])
                
                # Create a new slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Add header bar
                add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.8, 
                         fill_color=COLORS["primary"])
                
                # Add title
                add_text_box(slide, "Facilitation Notes Summary (Continued)", 0.5, 0.1, 9.0, 0.6, 
                           font_size=32, bold=True, color=COLORS["text_light"])
                
                # Create a container for the notes
                add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
                         0.5, 1.0, 9.0, FOOTER_Y - 1.2, 
                         fill_color=COLORS["light_alt"], 
                         line_color=COLORS["primary_light"],
                         line_width=1, shadow=True)
                
                # Reset y_position
                y_position = 1.2
    
    # Add footer
    # Add footer
    slide_number = slide_count_offset + quiz_count + discussion_count + 1
    add_footer(slide, content.get("title", "Untitled Presentation"), slide_number, total_slides, THEME["footer_style"])
    # add_footer(slide, content.get("title", "Untitled Presentation"), 
    #          total_slides, total_slides + 1, THEME["footer_style"])
    
    return slide

def create_closing_slide(prs, content, total_slides, slide_number):
    """Create an enhanced closing slide with visual elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background if enabled
    if THEME["use_gradients"]:
        add_gradient_background(prs, slide, COLORS["gradient_start"], COLORS["gradient_end"], angle=135)
    else:
        # Add solid color background
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT,
                 fill_color=COLORS["primary_dark"])
    
    # Add decorative corner accents if enabled
    if THEME["corner_accent"]:
        add_corner_accent(slide, COLORS["accent1"], 2.0, "top-right")
        add_corner_accent(slide, COLORS["accent2"], 1.5, "bottom-left")
    
    # Add title with enhanced styling
    title = "Thank You!"
    title_box = add_text_box(slide, title, 0.5, 1.5, 9, 1.5, font_size=48,
                           bold=True, color=COLORS["text_light"],
                           alignment=PP_ALIGN.CENTER, shadow=True)
    
    # Add subtitle with course title
    presentation_title = content.get("title", "Untitled Presentation")
    subtitle = f"Presentation: {presentation_title}"
    
    # Add a decorative line between title and subtitle
    line_y = 3.2
    add_shape(slide, MSO_SHAPE.RECTANGLE, 3.5, line_y, 3.0, 0.02,
             fill_color=COLORS["accent2"])
    
    # Add subtitle text
    subtitle_box = add_text_box(slide, subtitle, 0.5, line_y + 0.2, 9, 0.5,
                              font_size=28, italic=True, color=COLORS["text_light"],
                              alignment=PP_ALIGN.CENTER)
    
    # Add footer
    add_footer(slide, content.get("title", "Untitled Presentation"), slide_number, total_slides, THEME["footer_style"])
    
    return slide

def calculate_total_slides(content):
    """Calculate the total number of slides based on the content."""
    total = 0
    
    # Title slide
    total += 1
    
    # Agenda slides
    agenda_items = []
    agenda_items.append({"title": "Introduction", "items": [
        "Learning Outcomes",
        "Key Terms & Concepts"
    ]})
    
    content_slides = []
    for slide_content in content.get("slides", []):
        if slide_content.get("title", ""):
            title = clean_slide_title(slide_content.get("title", ""))
            if title:
                content_slides.append(title)
    
    if content_slides:
        agenda_items.append({"title": "Main Content", "items": content_slides})
    
    activities = []
    for idx, activity in enumerate(content.get("activities", [])):
        activity_title = activity.get("title", "")
        activities.append(activity_title)
    
    if activities:
        agenda_items.append({"title": "Activities", "items": activities})
    
    knowledge_items = []
    quiz_count = 0
    discussion_count = 0
    
    for idea in content.get("assessmentIdeas", []):
        idea_type = idea.get("type", "").lower()
        if "quiz" in idea_type:
            quiz_count += len([q for q in idea.get("exampleQuestions", []) if q.get("options")])
        elif "discussion" in idea_type:
            discussion_count += len(idea.get("exampleQuestions", []) if idea.get("exampleQuestions") else [])
    
    if quiz_count > 0:
        knowledge_items.append(f"Quiz Questions ({quiz_count})")
    
    if discussion_count > 0:
        knowledge_items.append(f"Discussion Questions ({discussion_count})")
    
    if knowledge_items:
        agenda_items.append({"title": "Test Your Knowledge", "items": knowledge_items})
    
    if content.get("furtherReadings", []):
        agenda_items.append({"title": "Additional Resources", "items": ["Further Readings & Resources"]})
    
    section_height = 0.5
    item_height = 0.35
    
    total_height_needed = 0
    for section in agenda_items:
        total_height_needed += section_height
        total_height_needed += len(section["items"]) * item_height
    
    available_height = FOOTER_Y - 1.2
    
    slides_needed = math.ceil(total_height_needed / available_height)
    total += slides_needed
    
    # Learning outcomes slide
    total += 1
    
    # Key terms slides
    key_terms = content.get("keyTerms", [])
    key_terms_per_slide = 4
    if key_terms:
        total += (len(key_terms) + key_terms_per_slide - 1) // key_terms_per_slide
    
    # Content slides
    total += len(content.get("slides", []))
    
    # Activity slides (2 per activity)
    total += len(content.get("activities", [])) * 2

    # Assessment slides
    for idea in content.get("assessmentIdeas", []):
        if "quiz" in idea.get("type", "").lower():
            total += len([q for q in idea.get("exampleQuestions", []) if q.get("options")]) * 2

    # Count discussion slides - now 2 slides per question (with guidance)
    for idea in content.get("assessmentIdeas", []):
        if "discussion" in idea.get("type", "").lower():
            total += len(idea.get("exampleQuestions", []) if idea.get("exampleQuestions") else []) * 2
    
    # Further Reading slides
    readings = content.get("furtherReadings", [])
    readings_per_slide = 2
    if readings:
        total += (len(readings) + readings_per_slide - 1) // readings_per_slide
    
    # Closing slide
    total += 1
    
    # Check if we need a facilitation notes slide
    has_facilitation_notes = False
    for activity in content.get("activities", []):
        description = activity.get("description", "")
        for pattern in ["Facilitation notes:", "Facilitation Notes:", "Facilitator notes:"]:
            if pattern in description:
                has_facilitation_notes = True
                break
        if has_facilitation_notes:
            break

    # Add facilitation notes slide if needed
    if has_facilitation_notes:
        total += 1
    
    return total



# Update the main function to handle multiple agenda slides
def main():
    """Main function to generate PowerPoint presentation."""
    # Define base directory for all content and output files
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

    # Define allowed content paths with fixed allowed names
    ALLOWED_CONTENT_PATHS = {
        "default": os.path.join(BASE_DIR, "content", "default_content.json"),
    }

    # Define allowed output directory
    OUTPUT_DIR = os.path.join(BASE_DIR, "output")
    if not os.path.exists(OUTPUT_DIR):
        try:
            os.makedirs(OUTPUT_DIR)
        except OSError as e:
            print(f"Error creating output directory: {e}")
            sys.exit(1)

    if len(sys.argv) != 3:
        print("Usage: python generate_pptx.py <content_key> <output_name>")
        sys.exit(1)

    content_key = sys.argv[1]
    output_name = sys.argv[2]

    # Strictly validate output name (only allow alphanumeric, underscore, hyphen)
    if not re.match(r'^[a-zA-Z0-9_\-]+$', output_name):
        print("Invalid output name. Only alphanumeric, underscore, hyphen allowed.")
        sys.exit(1)

    # Look up content path in the allowed paths dictionary
    if content_key == "default":
        content_path = ALLOWED_CONTENT_PATHS["default"]
    else:
        sys.exit(1)

    # Ensure content path exists
    if not os.path.exists(content_path):
        print(f"Content file does not exist: {content_path}")
        sys.exit(1)
    
    # Normalize paths to absolute paths to prevent path traversal
    content_path = os.path.abspath(content_path)

    # Verify content path is within allowed directories (defense in depth)
    if not any(content_path.startswith(os.path.abspath(allowed_path)) 
        for allowed_path in [ALLOWED_CONTENT_PATHS["default"]]):
            print("Security violation: Content path outside of allowed directory")
            sys.exit(1)

    # Construct output path with sanitized filename
    if not output_name.endswith('.pptx'):
        output_name += '.pptx'

    output_path = os.path.join(OUTPUT_DIR, output_name)
    
    # Normalize output path
    output_path = os.path.abspath(output_path)
    
    # Verify output path is within OUTPUT_DIR
    if not output_path.startswith(os.path.abspath(OUTPUT_DIR)):
        print("Security violation: Output path outside of allowed directory")
        sys.exit(1)

    # Load content from JSON file
    try:
        with open(content_path, 'r') as f:
            content = json.load(f)
    except json.JSONDecodeError:
        print("Invalid JSON format in content file")
        sys.exit(1)
    except Exception as e:
        print(f"Error loading content file: {e}")
        sys.exit(1)

    
    # Calculate total number of slides for internal tracking
    total_slides = calculate_total_slides(content)
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)
    
    # Create title slide (no page number)
    title_slide = create_title_slide(prs, content)
    
    # Create agenda slide (page 2)
    agenda_slides = create_agenda_slide(prs, content, total_slides)
    
    # Create learning outcomes slide (page 1)
    learning_outcomes_slide = create_learning_outcomes_slide(prs, content, total_slides)
    
    # Key terms slides (starting from page 2)
    key_terms_slides = create_key_terms_slide(prs, content, total_slides)
    
    # Calculate the starting slide number for content slides
    content_start_num = 2 + len(key_terms_slides)
    
    # Create content slides
    content_slides = create_content_slides(prs, content, total_slides)
    
    # Calculate the starting slide number for activity slides
    activity_start_num = content_start_num + len(content_slides)
    
    # Create activity slides
    activity_slides = create_activity_slides(prs, content, total_slides)
    
    # Calculate the starting slide number for quiz slides
    quiz_start_num = activity_start_num + len(activity_slides)
    
    # Create quiz slides
    quiz_slides = create_quiz_slides(prs, content, total_slides)
    
    # Count quiz slides
    quiz_count = len(quiz_slides)
    
    # Calculate the starting slide number for discussion slides
    discussion_start_num = quiz_start_num + quiz_count
    
    # Create discussion slides with facilitator guidance
    discussion_slides = create_discussion_slides(prs, content, total_slides)
    
    # Count discussion slides - now only 1 per question
    discussion_count = len(discussion_slides)
    
    # Calculate the starting slide number for further readings
    readings_start_num = discussion_start_num + discussion_count
    
    # Create further readings slides with version 7 styling
    readings_slides = create_further_readings_slides(prs, content, total_slides)
    
    # Count readings slides
    readings_count = len(readings_slides)
    
    # Calculate the final slide number
    closing_num = readings_start_num + readings_count
    
    # Create facilitation notes summary slide if applicable
    facilitation_slide = create_facilitation_notes_slide(prs, content, total_slides)
    if facilitation_slide:
        # Increment total_slides since we added a new slide
        total_slides += 1
        # Update closing_num
        closing_num += 1
    
    # Create closing slide with updated page number
    closing_slide = create_closing_slide(prs, content, total_slides, closing_num)
    
    # Save presentation
    prs.save(output_path)
    print(f"PowerPoint presentation saved to {output_path}")

if __name__ == "__main__":
    main()

def create_pptx(content: dict, output_path: str):
    """
    Generate a PowerPoint presentation based on the provided content.

    Args:
        content (dict): The content for the presentation. Expected keys include:
                        - title (str): Title of the presentation.
                        - slides (list): List of slides, where each slide is a dict with keys:
                            - heading (str): Slide heading.
                            - body (str): Slide body text.
        output_path (str): The file path to save the generated PPTX file.
    """
    try:
        # Validate output path
        BASE_DIR = os.path.dirname(os.path.abspath(__file__))
        
        # Normalize output path to prevent path traversal attacks
        normalized_output_path = os.path.abspath(output_path)
        
        # Verify the output path is within one of the allowed directories
        # Get parent directory to check if it's in a system temp directory or the application's output directory
        output_parent = os.path.dirname(normalized_output_path)
        output_dir = os.path.join(BASE_DIR, "output")
        
        # Check if output is in system temp dir or the output dir
        is_in_temp = output_parent.startswith(os.path.abspath(tempfile.gettempdir()))
        is_in_output = normalized_output_path.startswith(os.path.abspath(output_dir))
        
        if not (is_in_temp or is_in_output):
            raise ValueError("Security violation: Output path must be in allowed directories")
            
        # Create a new PowerPoint presentation
        prs = Presentation()

        # Set slide dimensions to 16:9 aspect ratio
        prs.slide_width = Inches(SLIDE_WIDTH)
        prs.slide_height = Inches(SLIDE_HEIGHT)

        # Calculate total number of slides for internal tracking
        total_slides = calculate_total_slides(content)

        # Create title slide
        create_title_slide(prs, content)

        # Create agenda slides
        create_agenda_slide(prs, content, total_slides)

        # Create learning outcomes slide
        create_learning_outcomes_slide(prs, content, total_slides)

        # Create key terms slides
        create_key_terms_slide(prs, content, total_slides)

        # Create content slides
        create_content_slides(prs, content, total_slides)

        # Create activity slides
        create_activity_slides(prs, content, total_slides)

        # Create quiz slides
        create_quiz_slides(prs, content, total_slides)

        # Create discussion slides
        create_discussion_slides(prs, content, total_slides)

        # Create further readings slides
        create_further_readings_slides(prs, content, total_slides)

        # Create facilitation notes slide if applicable
        facilitation_slide = create_facilitation_notes_slide(prs, content, total_slides)
        if facilitation_slide:
            total_slides += 1  # Increment total slides if facilitation notes slide is added

        # Create closing slide
        create_closing_slide(prs, content, total_slides, total_slides)

        # Save the presentation to the specified output path
        prs.save(normalized_output_path)
    except Exception as e:
        print(f"Error generating PPTX: {e}")
        raise