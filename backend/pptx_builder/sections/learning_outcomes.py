import re
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from ..constants import COLORS, SLIDE_WIDTH, THEME, FOOTER_Y
from ..shapes import add_shape, add_text_box, add_footer
from ..localization import t


def create_learning_outcomes_slide(prs, content, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.8, fill_color=COLORS["primary"]
    )
    add_text_box(
        slide,
        t("learningOutcomes"),
        0.5,
        0.1,
        9,
        0.6,
        font_size=36,
        bold=True,
        color=COLORS["text_light"],
    )
    ct_names = t("contentTypeNames")
    content_type = content.get("contentType", "lecture")
    content_type_display = (
        ct_names.get(content_type, "lecture")
        if isinstance(ct_names, dict)
        else "lecture"
    )
    intro_text = t("byTheEnd", contentType=content_type_display)
    add_text_box(
        slide,
        intro_text,
        0.5,
        1.0,
        9,
        0.5,
        font_size=20,
        italic=True,
        color=COLORS["dark"],
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        0.3,
        1.7,
        9.4,
        3.0,
        fill_color=COLORS["light"],
        opacity=0.9,
        line_color=COLORS["primary_light"],
        line_width=1,
    )
    learning_outcomes = content.get("learningOutcomes", [])
    y = 2.0
    bullet_colors = [COLORS["emerald"], COLORS["medium_purple"], COLORS["emerald"]]
    for idx, outcome in enumerate(learning_outcomes):
        cleaned = re.sub(r"^\d+\.\s*", "", outcome)
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            0.7,
            y,
            0.15,
            0.15,
            fill_color=bullet_colors[idx % len(bullet_colors)],
        )
        add_text_box(
            slide,
            cleaned,
            1.0,
            y - 0.125,
            8.5,
            0.4,
            font_size=20,
            color=COLORS["text"],
            vertical_alignment=MSO_ANCHOR.MIDDLE,
        )
        y += 0.6
    add_footer(
        slide,
        (content.get("title") or t("untitledPresentation")),
        1,
        total_slides,
        THEME["footer_style"],
    )
    return slide
