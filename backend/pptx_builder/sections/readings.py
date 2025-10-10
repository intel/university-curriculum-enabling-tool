from pptx.enum.shapes import MSO_SHAPE
from ..constants import COLORS, SLIDE_WIDTH, FOOTER_Y, THEME
from ..shapes import add_shape, add_text_box, add_footer
from ..localization import t


def create_further_readings_slides(prs, content, total_slides):
    readings = content.get("furtherReadings", [])
    if not readings:
        return []
    slides = []
    slide_count_offset = (
        2
        + len(content.get("keyTerms", [])) // 4
        - 1
        + len([s for i, s in enumerate(content.get("slides", [])) if i != 1])
        + len(content.get("activities", [])) * 2
        + 1
    )
    quiz_count = 0
    for idea in content.get("assessmentIdeas", []):
        if "quiz" in idea.get("type", "").lower():
            quiz_count += (
                len([q for q in idea.get("exampleQuestions", []) if q.get("options")])
                * 2
            )
    discussion_count = 0
    for idea in content.get("assessmentIdeas", []):
        if "discussion" in idea.get("type", "").lower():
            discussion_count += (
                len(
                    idea.get("exampleQuestions", [])
                    if idea.get("exampleQuestions")
                    else []
                )
                * 2
            )
    readings_per_slide = 2
    total_readings = len(readings)
    slides_needed = (total_readings + readings_per_slide - 1) // readings_per_slide
    for slide_idx in range(slides_needed):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slides.append(slide)
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            SLIDE_WIDTH,
            0.8,
            fill_color=COLORS["primary"],
        )
        title = t("furtherReadings")
        if slide_idx > 0:
            title += t("continued")
        add_text_box(
            slide,
            title,
            0.5,
            0.1,
            9,
            0.6,
            font_size=32,
            bold=True,
            color=COLORS["text_light"],
        )
        if THEME["content_box_shadow"]:
            add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                0.3,
                1.0,
                9.4,
                FOOTER_Y - 1.2,
                fill_color=COLORS["light_alt"],
                opacity=0.7,
                line_color=COLORS["primary_light"],
                line_width=1,
                shadow=True,
            )
        start_idx = slide_idx * readings_per_slide
        end_idx = min(start_idx + readings_per_slide, total_readings)
        readings_for_slide = readings[start_idx:end_idx]
        y = 1.2
        for i, reading in enumerate(readings_for_slide):
            reading_title = reading.get("title") or t("untitledReading")
            add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                0.7,
                y,
                0.1,
                0.4,
                fill_color=COLORS["primary"],
            )
            add_text_box(
                slide,
                reading_title,
                0.9,
                y,
                8.3,
                0.4,
                font_size=20,
                bold=True,
                color=COLORS["primary"],
            )
            y += 0.5
            reading_author = reading.get("author") or t("unknownAuthor")
            reading_description = reading.get("readingDescription", "")
            add_text_box(
                slide,
                f"{t('author')} {reading_author}",
                0.9,
                y,
                8.3,
                0.3,
                font_size=16,
                italic=True,
                color=COLORS["primary"],
            )
            y += 0.4
            add_text_box(
                slide,
                reading_description,
                0.9,
                y,
                8.3,
                0.6,
                font_size=16,
                color=COLORS["text"],
            )
            if i < len(readings_for_slide) - 1:
                y += 0.8
                add_shape(
                    slide,
                    MSO_SHAPE.RECTANGLE,
                    0.7,
                    y,
                    8.5,
                    0.01,
                    fill_color=COLORS["primary_light"],
                    opacity=0.5,
                )
                y += 0.2
            else:
                y += 0.8
        slide_number = (
            slide_count_offset + quiz_count + discussion_count + slide_idx + 1
        )
        add_footer(
            slide,
            (content.get("title") or t("untitledPresentation")),
            slide_number,
            total_slides,
            THEME["footer_style"],
        )
    return slides
