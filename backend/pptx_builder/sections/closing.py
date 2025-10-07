from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from ..constants import COLORS, SLIDE_WIDTH, SLIDE_HEIGHT, THEME
from ..shapes import (
    add_text_box,
    add_shape,
    add_corner_accent,
    add_gradient_background,
    add_footer,
)
from ..localization import t


def create_closing_slide(prs, content, total_slides, slide_number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if THEME["use_gradients"]:
        add_gradient_background(
            prs, slide, COLORS["gradient_start"], COLORS["gradient_end"], angle=135
        )
    else:
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            SLIDE_WIDTH,
            SLIDE_HEIGHT,
            fill_color=COLORS["primary_dark"],
        )
    if THEME["corner_accent"]:
        add_corner_accent(slide, COLORS["accent1"], 2.0, "top-right")
        add_corner_accent(slide, COLORS["accent2"], 1.5, "bottom-left")
    title = t("thankYou")
    add_text_box(
        slide,
        title,
        0.5,
        1.5,
        9,
        1.5,
        font_size=48,
        bold=True,
        color=COLORS["text_light"],
        alignment=PP_ALIGN.CENTER,
        shadow=True,
    )
    presentation_title = content.get("title") or t("untitledPresentation")
    subtitle = f"{t('presentation')} {presentation_title}"
    line_y = 3.2
    add_shape(
        slide, MSO_SHAPE.RECTANGLE, 3.5, line_y, 3.0, 0.02, fill_color=COLORS["accent2"]
    )
    add_text_box(
        slide,
        subtitle,
        0.5,
        line_y + 0.2,
        9,
        0.5,
        font_size=28,
        italic=True,
        color=COLORS["text_light"],
        alignment=PP_ALIGN.CENTER,
    )
    add_footer(
        slide,
        (content.get("title") or t("untitledPresentation")),
        slide_number,
        total_slides,
        THEME["footer_style"],
    )
    return slide
