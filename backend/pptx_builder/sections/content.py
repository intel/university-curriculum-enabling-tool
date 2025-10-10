import json
from pptx.enum.shapes import MSO_SHAPE
from ..constants import (
    COLORS,
    THEME,
    SLIDE_WIDTH,
    FOOTER_Y,
    CONTENT_START_Y,
    MAIN_BULLET_INDENT,
)
from ..shapes import (
    add_gradient_background,
    add_shape,
    add_corner_accent,
    add_text_box,
    add_footer,
)
from ..localization import t
from ..utils import clean_slide_title
from pptx.util import Inches, Pt


def create_content_slides(prs, content, total_slides):
    slides_data = content.get("slides", [])
    if not slides_data:
        return []
    result = []
    slide_count_offset = 2 + len(content.get("keyTerms", [])) // 4
    for slide_idx, slide_content in enumerate(slides_data):
        if slide_idx == 1:  # skip second slide as per original logic
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        result.append(slide)
        if THEME["use_gradients"]:
            add_gradient_background(
                prs, slide, COLORS["primary"], COLORS["primary_dark"], angle=0
            )
            add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                0,
                0.8,
                SLIDE_WIDTH,
                (5.625 - 0.8),
                fill_color=COLORS["background"],
                opacity=0.9,
            )
        else:
            add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                SLIDE_WIDTH,
                0.8,
                fill_color=COLORS["royal_blue"],
            )
        if THEME["corner_accent"]:
            accent_color = [COLORS["accent1"], COLORS["accent2"], COLORS["accent3"]][
                slide_idx % 3
            ]
            add_corner_accent(slide, accent_color, 1.0, "bottom-right")
        cleaned_title = clean_slide_title(slide_content.get("title", ""))
        add_text_box(
            slide,
            cleaned_title,
            0.5,
            0.1,
            9,
            0.6,
            font_size=32,
            bold=True,
            color=COLORS["text_light"],
        )
        points = slide_content.get("content", [])
        if THEME["content_box_shadow"]:
            content_height = FOOTER_Y - CONTENT_START_Y - 0.2
            add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                0.3,
                CONTENT_START_Y - 0.1,
                9.4,
                content_height,
                fill_color=COLORS["light_alt"],
                opacity=0.7,
                line_color=COLORS["primary_light"],
                line_width=1,
                shadow=True,
            )
        tb = slide.shapes.add_textbox(
            Inches(MAIN_BULLET_INDENT),
            Inches(CONTENT_START_Y),
            Inches(9 - MAIN_BULLET_INDENT),
            Inches(FOOTER_Y - CONTENT_START_Y - 0.3),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        has_sub = any(
            p.strip().startswith(("  ", "\\t", "-"))
            for p in points
            if isinstance(p, str)
        )
        p = tf.paragraphs[0]
        first = True
        for point in points:
            text = point if isinstance(point, str) else json.dumps(point)
            is_bullet = False
            level = 0
            if text.strip().startswith(("•", "*")):
                is_bullet = True
                text = text.strip()[1:].strip()
            elif text.strip().startswith("-"):
                is_bullet = True
                level = 1
                text = text.strip()[1:].strip()
            elif text.strip().startswith("  ") or text.strip().startswith("\\t"):
                is_bullet = True
                level = 1
                text = text.strip()
            elif not has_sub:
                is_bullet = True
            if not first:
                p = tf.add_paragraph()
            else:
                first = False
            if is_bullet:
                p.level = level
                try:
                    p.bullet.visible = True
                except:
                    if THEME["modern_bullets"]:
                        text = ("◦ " if level > 0 else "• ") + text
            run = p.add_run()
            run.text = text
            font = run.font
            font.size = Pt(18) if level == 0 else Pt(16)
            font.bold = not is_bullet and level == 0
            font.color.rgb = COLORS["text"]
        notes = slide_content.get("notes", "")
        if notes:
            if not slide.has_notes_slide:
                slide.notes_slide
            slide.notes_slide.notes_text_frame.text = (
                notes if isinstance(notes, str) else json.dumps(notes)
            )
        adjusted_idx = slide_idx if slide_idx < 1 else slide_idx - 1
        slide_number = slide_count_offset + adjusted_idx + 1
        add_footer(
            slide,
            (content.get("title") or t("untitledPresentation")),
            slide_number,
            total_slides,
            THEME["footer_style"],
        )
    return result
