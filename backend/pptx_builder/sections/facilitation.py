from pptx.enum.shapes import MSO_SHAPE
from ..constants import COLORS, SLIDE_WIDTH, FOOTER_Y, THEME, GLOBAL_LANG, LABELS
from ..shapes import add_shape, add_text_box, add_footer
from ..localization import t
from ..utils import extract_facilitation_content


def create_facilitation_notes_slide(prs, content, total_slides):
    activities = content.get("activities", [])
    slide_count_offset = (
        2
        + len(content.get("keyTerms", [])) // 4
        - 1
        + len([s for i, s in enumerate(content.get("slides", [])) if i != 1])
        + len(content.get("activities", [])) * 2
        + 1
        + len(content.get("furtherReadings", [])) // 2
    )
    quiz_count = 0
    for idea in content.get("assessmentIdeas", []):
        if "quiz" in idea.get("type", "").lower():
            quiz_count += (
                len([q for q in idea.get("exampleQuestions", []) if q.get("options")])
                * 2
            )
    discussion_count = 0
    for idea in content.get("assessmentIdeas", []):
        if "discussion" in idea.get("type", "").lower():
            discussion_count += (
                len(
                    idea.get("exampleQuestions", [])
                    if idea.get("exampleQuestions")
                    else []
                )
                * 2
            )
    has_notes = False
    for activity in activities:
        _, facilitation_notes, _ = extract_facilitation_content(
            activity.get("description", "")
        )
        if facilitation_notes:
            has_notes = True
            break
    if not has_notes:
        return None
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.8, fill_color=COLORS["primary"]
    )
    add_text_box(
        slide,
        LABELS[GLOBAL_LANG].get(
            "facilitationNotesSummary", "Facilitation Notes Summary"
        ),
        0.5,
        0.1,
        9.0,
        0.6,
        font_size=32,
        bold=True,
        color=COLORS["text_light"],
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        0.5,
        1.0,
        9.0,
        FOOTER_Y - 1.2,
        fill_color=COLORS["light_alt"],
        line_color=COLORS["primary_light"],
        line_width=1,
        shadow=True,
    )
    y = 1.2
    for idx, activity in enumerate(activities):
        title = activity.get("title", "")
        description = activity.get("description", "")
        _, facilitation_notes, _ = extract_facilitation_content(description)
        if facilitation_notes:
            add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                0.7,
                y,
                0.1,
                0.4,
                fill_color=COLORS["activity_green"],
            )
            add_text_box(
                slide,
                title,
                0.9,
                y,
                8.3,
                0.4,
                font_size=18,
                bold=True,
                color=COLORS["primary"],
            )
            y += 0.5
            notes_text = facilitation_notes.replace("Facilitation Notes: ", "")
            add_text_box(
                slide, notes_text, 0.9, y, 8.3, 0.6, font_size=14, color=COLORS["text"]
            )
            y += 0.8
            if idx < len(activities) - 1:
                add_shape(
                    slide,
                    MSO_SHAPE.RECTANGLE,
                    0.7,
                    y,
                    8.5,
                    0.01,
                    fill_color=COLORS["primary_light"],
                    opacity=0.5,
                )
                y += 0.3
            if y > FOOTER_Y - 0.5:
                add_text_box(
                    slide,
                    LABELS[GLOBAL_LANG].get(
                        "continuedNextSlide", "Continued on next slide..."
                    ),
                    0.9,
                    y - 0.3,
                    8.3,
                    0.3,
                    font_size=12,
                    italic=True,
                    color=COLORS["text_muted"],
                )
                add_footer(
                    slide,
                    (content.get("title") or t("untitledPresentation")),
                    total_slides - 1,
                    total_slides + 1,
                    THEME["footer_style"],
                )
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_shape(
                    slide,
                    MSO_SHAPE.RECTANGLE,
                    0,
                    0,
                    SLIDE_WIDTH,
                    0.8,
                    fill_color=COLORS["primary"],
                )
                add_text_box(
                    slide,
                    LABELS[GLOBAL_LANG].get(
                        "facilitationNotesSummary", "Facilitation Notes Summary"
                    )
                    + LABELS[GLOBAL_LANG].get("continued", " (continued)"),
                    0.5,
                    0.1,
                    9.0,
                    0.6,
                    font_size=32,
                    bold=True,
                    color=COLORS["text_light"],
                )
                add_shape(
                    slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    0.5,
                    1.0,
                    9.0,
                    FOOTER_Y - 1.2,
                    fill_color=COLORS["light_alt"],
                    line_color=COLORS["primary_light"],
                    line_width=1,
                    shadow=True,
                )
                y = 1.2
    slide_number = slide_count_offset + quiz_count + discussion_count + 1
    add_footer(
        slide,
        (content.get("title") or t("untitledPresentation")),
        slide_number,
        total_slides,
        THEME["footer_style"],
    )
    return slide
