import math
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from ..constants import COLORS, SLIDE_WIDTH, FOOTER_Y, THEME
from ..shapes import add_shape, add_text_box, add_footer
from ..localization import t
from ..utils import clean_slide_title


def create_agenda_slide(prs, content, total_slides):
    agenda_items = []
    intro_items = [t("learningOutcomes"), t("keyTerms")]
    agenda_items.append({"title": t("introduction"), "items": intro_items})
    content_slides = []
    for slide_content in content.get("slides", []):
        title = clean_slide_title(slide_content.get("title", ""))
        if title:
            content_slides.append(title)
    if content_slides:
        agenda_items.append({"title": t("mainContent"), "items": content_slides})
    activities = [a.get("title", "") for a in content.get("activities", [])]
    if activities:
        agenda_items.append({"title": t("activities"), "items": activities})
    knowledge_items = []
    quiz_count = 0
    discussion_count = 0
    for idea in content.get("assessmentIdeas", []):
        idea_type = idea.get("type", "").lower()
        if "quiz" in idea_type:
            quiz_count += len(
                [q for q in idea.get("exampleQuestions", []) if q.get("options")]
            )
        elif "discussion" in idea_type:
            discussion_count += len(
                idea.get("exampleQuestions", []) if idea.get("exampleQuestions") else []
            )
    if quiz_count > 0:
        knowledge_items.append(t("quizQuestions", count=quiz_count))
    if discussion_count > 0:
        knowledge_items.append(t("discussionQuestions", count=discussion_count))
    if knowledge_items:
        agenda_items.append({"title": t("testYourKnowledge"), "items": knowledge_items})
    if content.get("furtherReadings", []):
        agenda_items.append(
            {"title": t("additionalResources"), "items": [t("furtherReadings")]}
        )
    section_height = 0.5
    item_height = 0.35
    total_height_needed = 0
    for section in agenda_items:
        total_height_needed += section_height + len(section["items"]) * item_height
    available_height = FOOTER_Y - 1.2
    slides_needed = math.ceil(total_height_needed / available_height)
    agenda_slides = []
    section_start_idx = 0
    item_start_idx = 0
    consumed = []  # track per slide
    for slide_idx in range(slides_needed):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        agenda_slides.append(slide)
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            SLIDE_WIDTH,
            0.8,
            fill_color=COLORS["primary"],
        )
        title = t("agenda")
        if slide_idx > 0:
            title += t("agendaContinued", idx=slide_idx + 1, total=slides_needed)
        add_text_box(
            slide,
            title,
            0.5,
            0.1,
            9,
            0.6,
            font_size=36,
            bold=True,
            color=COLORS["text_light"],
        )
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            0.3,
            0.9,
            9.4,
            FOOTER_Y - 1.1,
            fill_color=COLORS["light"],
            opacity=0.9,
            line_color=COLORS["primary_light"],
            line_width=1,
        )
        y = 1.1
        max_y = FOOTER_Y - 0.3
        current_section_idx = section_start_idx
        current_item_start = item_start_idx
        while current_section_idx < len(agenda_items) and y < max_y:
            section = agenda_items[current_section_idx]
            if y + section_height > max_y:
                break
            add_text_box(
                slide,
                section["title"],
                0.7,
                y,
                8.5,
                section_height,
                font_size=24,
                bold=True,
                color=COLORS["primary"],
            )
            y += section_height
            items = section["items"][current_item_start:]
            for idx, item in enumerate(items):
                if y + item_height > max_y:
                    item_start_idx = current_item_start + idx
                    break
                tb = slide.shapes.add_textbox(
                    Inches(1.0), Inches(y), Inches(8.0), Inches(item_height)
                )
                p = tb.text_frame.paragraphs[0]
                p.level = 1
                try:
                    p.bullet.visible = True
                except:
                    item = f"• {item}"
                run = p.add_run()
                run.text = item
                run.font.size = Pt(18)
                run.font.color.rgb = COLORS["text"]
                y += item_height
            else:
                current_section_idx += 1
                current_item_start = 0
                item_start_idx = 0
                section_start_idx = current_section_idx
                continue
            break
        add_footer(
            slide,
            (content.get("title") or t("untitledPresentation")),
            slide_idx + 2,
            total_slides,
            THEME["footer_style"],
        )
    return agenda_slides
