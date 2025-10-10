from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from ..constants import COLORS, SLIDE_WIDTH, SLIDE_HEIGHT, THEME
from ..shapes import add_text_box, add_shape, add_corner_accent, add_gradient_background
from ..localization import t


def create_title_slide(prs, content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if THEME["use_gradients"]:
        add_gradient_background(
            prs, slide, COLORS["gradient_start"], COLORS["gradient_end"], angle=135
        )
    else:
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            SLIDE_WIDTH,
            SLIDE_HEIGHT,
            fill_color=COLORS["primary_dark"],
        )
    if THEME["corner_accent"]:
        add_corner_accent(slide, COLORS["accent1"], 2.0, "top-right")
        add_corner_accent(slide, COLORS["accent2"], 1.5, "bottom-left")
    title = content.get("title", t("untitledPresentation"))
    add_text_box(
        slide,
        title,
        0.5,
        1.5,
        9,
        1.5,
        font_size=48,
        bold=True,
        color=COLORS["text_light"],
        alignment=PP_ALIGN.CENTER,
        shadow=True,
    )
    ct_names = t("contentTypeNames")
    diff_names = t("difficultyNames")
    content_type = content.get("contentType", "lecture")
    difficulty = content.get("difficultyLevel", "intermediate")
    ct_disp = (
        ct_names.get(content_type, "Lecture")
        if isinstance(ct_names, dict)
        else "Lecture"
    )
    diff_disp = (
        diff_names.get(difficulty, "Intermediate Level")
        if isinstance(diff_names, dict)
        else "Intermediate Level"
    )
    subtitle = f"{ct_disp} | {diff_disp}"
    line_y = 3.2
    add_shape(
        slide, MSO_SHAPE.RECTANGLE, 3.5, line_y, 3.0, 0.02, fill_color=COLORS["accent2"]
    )
    add_text_box(
        slide,
        subtitle,
        0.5,
        line_y + 0.2,
        9,
        0.5,
        font_size=28,
        italic=True,
        color=COLORS["text_light"],
        alignment=PP_ALIGN.CENTER,
    )
    return slide
