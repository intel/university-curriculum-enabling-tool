import json
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from ..constants import COLORS, SLIDE_WIDTH, SLIDE_HEIGHT, FOOTER_Y, THEME
from ..shapes import add_shape, add_text_box, add_footer
from ..localization import t
from ..utils import clean_activity_title, extract_facilitation_content


def create_activity_slides(prs, content, total_slides):
    activities = content.get("activities", [])
    if not activities:
        return []
    slides = []
    slide_count_offset = (
        2
        + len(content.get("keyTerms", [])) // 4
        - 1
        + len([s for i, s in enumerate(content.get("slides", [])) if i != 1])
        + 1
    )
    for act_idx, activity in enumerate(activities):
        original_title = activity.get("title", "") or t("untitledActivity")
        clean_title = clean_activity_title(original_title)
        main_slide = prs.slides.add_slide(prs.slide_layouts[6])
        slides.append(main_slide)
        materials_slide = prs.slides.add_slide(prs.slide_layouts[6])
        slides.append(materials_slide)
        # Header bars
        for slide in (main_slide, materials_slide):
            add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                SLIDE_WIDTH,
                0.8,
                fill_color=COLORS["activity_blue"],
            )
        add_text_box(
            main_slide,
            t("activity", num=act_idx + 1, title=clean_title),
            0.5,
            0.07,
            9.0,
            0.6,
            font_size=22,
            bold=True,
            color=COLORS["text_light"],
        )
        add_text_box(
            materials_slide,
            f"{t('activity', num=act_idx + 1, title=clean_title)} {t('materialsSuffix')}",
            0.5,
            0.07,
            9.0,
            0.6,
            font_size=22,
            bold=True,
            color=COLORS["text_light"],
        )
        # Badge
        for slide in (main_slide, materials_slide):
            add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                0.5,
                0.9,
                5.0,
                0.5,
                fill_color=COLORS["activity_purple"],
            )
        activity_type = activity.get("type", "Exercise")
        activity_duration = activity.get("duration", "20 minutes")
        type_duration_text = (
            f"{t('type')}: {activity_type} | {t('duration')}: {activity_duration}"
        )
        for slide in (main_slide, materials_slide):
            add_text_box(
                slide,
                type_duration_text,
                0.6,
                0.9,
                4.8,
                0.5,
                font_size=16,
                italic=True,
                color=COLORS["text_light"],
                vertical_alignment=MSO_ANCHOR.MIDDLE,
            )
        # Containers
        for slide in (main_slide, materials_slide):
            add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                0.5,
                1.5,
                9.0,
                3.5,
                fill_color=COLORS["background"],
                line_color=COLORS["activity_purple"],
                line_width=1,
            )
        activity_description = activity.get("description", "")
        clean_description, facilitation_notes, learning_objectives = (
            extract_facilitation_content(activity_description)
        )
        for slide in (main_slide, materials_slide):
            add_text_box(
                slide,
                clean_description,
                0.7,
                1.7,
                8.6,
                0.6,
                font_size=20,
                color=COLORS["text"],
            )
        if facilitation_notes or learning_objectives:
            combined = []
            if facilitation_notes:
                combined.append(f"{t('facilitationNotesLabel')} {facilitation_notes}")
            if learning_objectives:
                combined.append(f"{t('learningObjectiveLabel')} {learning_objectives}")
            notes_text = "\n\n".join(combined)
            for slide in (main_slide, materials_slide):
                if not slide.has_notes_slide:
                    slide.notes_slide
                slide.notes_slide.notes_text_frame.text = notes_text
        if facilitation_notes:
            for slide in (main_slide, materials_slide):
                add_shape(
                    slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    8.5,
                    0.9,
                    1.0,
                    0.5,
                    fill_color=COLORS["activity_green"],
                )
                add_text_box(
                    slide,
                    t("notesAvailable"),
                    8.6,
                    0.95,
                    0.8,
                    0.4,
                    font_size=12,
                    bold=True,
                    color=COLORS["text_light"],
                    alignment=PP_ALIGN.CENTER,
                )
        # Instructions
        instructions_y = 2.4
        add_shape(
            main_slide,
            MSO_SHAPE.RECTANGLE,
            0.7,
            instructions_y + 0.1,
            0.1,
            2.0,
            fill_color=COLORS["activity_blue"],
        )
        add_text_box(
            main_slide,
            t("instructions"),
            0.9,
            instructions_y,
            8.3,
            0.4,
            font_size=22,
            bold=True,
            color=COLORS["text"],
        )
        instructions = activity.get("instructions", [])
        instr_tb = main_slide.shapes.add_textbox(
            Inches(0.9), Inches(instructions_y + 0.5), Inches(8.3), Inches(1.5)
        )
        frame = instr_tb.text_frame
        frame.word_wrap = True
        for idx, instruction in enumerate(instructions):
            p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            text = (
                instruction if isinstance(instruction, str) else json.dumps(instruction)
            )
            p.text = f"{idx + 1}. {text}"
            for run in p.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = COLORS["text"]
        # Materials
        materials_y = 2.4
        add_shape(
            materials_slide,
            MSO_SHAPE.RECTANGLE,
            0.7,
            materials_y + 0.1,
            0.1,
            2.0,
            fill_color=COLORS["activity_green"],
        )
        add_text_box(
            materials_slide,
            t("materialsNeeded"),
            0.9,
            materials_y,
            8.3,
            0.4,
            font_size=22,
            bold=True,
            color=COLORS["text"],
        )
        materials = activity.get(
            "materials", ["Gaudi-3 optimization tools", "Neural network models"]
        )
        mat_tb = materials_slide.shapes.add_textbox(
            Inches(0.9), Inches(materials_y + 0.5), Inches(8.3), Inches(1.5)
        )
        mat_frame = mat_tb.text_frame
        mat_frame.word_wrap = True
        for idx, material in enumerate(materials):
            p = mat_frame.paragraphs[0] if idx == 0 else mat_frame.add_paragraph()
            text = material if isinstance(material, str) else json.dumps(material)
            p.text = f"• {text}"
            for run in p.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = COLORS["text"]
        # Bottom accent triangles
        for slide in (main_slide, materials_slide):
            add_shape(
                slide,
                MSO_SHAPE.RIGHT_TRIANGLE,
                0,
                SLIDE_HEIGHT - 1.5,
                1.5,
                1.5,
                fill_color=COLORS["activity_orange"],
            )
            add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                0.5,
                FOOTER_Y - 0.05,
                9.0,
                0.01,
                fill_color=COLORS["primary_light"],
            )
        presentation_title = content.get("title") or t("untitledPresentation")
        main_num = slide_count_offset + (act_idx * 2) + 1
        materials_num = slide_count_offset + (act_idx * 2) + 2
        add_text_box(
            main_slide,
            presentation_title,
            0.5,
            FOOTER_Y,
            8.0,
            0.3,
            font_size=10,
            color=COLORS["primary"],
            italic=True,
        )
        add_text_box(
            main_slide,
            f"{main_num}",
            9.0,
            FOOTER_Y,
            0.5,
            0.3,
            font_size=10,
            color=COLORS["primary"],
            alignment=PP_ALIGN.RIGHT,
        )
        add_text_box(
            materials_slide,
            presentation_title,
            0.5,
            FOOTER_Y,
            8.0,
            0.3,
            font_size=10,
            color=COLORS["primary"],
            italic=True,
        )
        add_text_box(
            materials_slide,
            f"{materials_num}",
            9.0,
            FOOTER_Y,
            0.5,
            0.3,
            font_size=10,
            color=COLORS["primary"],
            alignment=PP_ALIGN.RIGHT,
        )
    return slides
