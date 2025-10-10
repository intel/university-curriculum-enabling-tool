import json
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from ..constants import COLORS, SLIDE_WIDTH, FOOTER_Y, THEME
from ..shapes import add_shape, add_text_box, add_footer
from ..localization import t


def create_quiz_slides(prs, content, total_slides):
    assessment_ideas = content.get("assessmentIdeas", [])
    slides = []
    slide_count_offset = (
        2
        + len(content.get("keyTerms", [])) // 4
        - 1
        + len([s for i, s in enumerate(content.get("slides", [])) if i != 1])
        + len(content.get("activities", [])) * 2
        + 1
    )
    quiz_slide_count = 0
    for idea in assessment_ideas:
        if "quiz" not in idea.get("type", "").lower():
            continue
        for q_idx, question in enumerate(idea.get("exampleQuestions", [])):
            question_text = question.get("question", "Example question")
            options = question.get("options", [])
            if not options:
                continue
            q_slide = prs.slides.add_slide(prs.slide_layouts[6])
            slides.append(q_slide)
            quiz_slide_count += 1
            add_shape(
                q_slide,
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                SLIDE_WIDTH,
                1.0,
                fill_color=COLORS["primary"],
            )
            add_text_box(
                q_slide,
                t("quizQuestion", num=q_idx + 1),
                0.5,
                0.2,
                9.0,
                0.6,
                font_size=36,
                bold=True,
                color=COLORS["text_light"],
                alignment=PP_ALIGN.CENTER,
                vertical_alignment=MSO_ANCHOR.MIDDLE,
            )
            add_shape(
                q_slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                0.5,
                1.1,
                9.0,
                0.8,
                fill_color=COLORS["light"],
                line_color=COLORS["light"],
                line_width=1,
            )
            add_text_box(
                q_slide,
                question_text,
                0.7,
                1.2,
                8.6,
                0.6,
                font_size=20,
                bold=True,
                color=COLORS["text"],
            )
            options_per_row = 2
            option_width = 4.3
            option_height = 1.0
            option_gap = 0.4
            start_y = 2.2
            for opt_idx, option in enumerate(options):
                row = opt_idx // options_per_row
                col = opt_idx % options_per_row
                ox = 0.5 + col * (option_width + option_gap)
                oy = start_y + row * (option_height + 0.4)
                add_shape(
                    q_slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    ox,
                    oy,
                    option_width,
                    option_height,
                    fill_color=COLORS["light"],
                    line_color=COLORS["light"],
                )
                circle_size = 0.6
                cx = ox + 0.2
                cy = oy + (option_height - circle_size) / 2
                add_shape(
                    q_slide,
                    MSO_SHAPE.OVAL,
                    cx,
                    cy,
                    circle_size,
                    circle_size,
                    fill_color=COLORS["primary"],
                )
                add_text_box(
                    q_slide,
                    chr(65 + opt_idx),
                    cx,
                    cy,
                    circle_size,
                    circle_size,
                    font_size=24,
                    bold=True,
                    color=COLORS["text_light"],
                    alignment=PP_ALIGN.CENTER,
                    vertical_alignment=MSO_ANCHOR.MIDDLE,
                )
                text_x = cx + circle_size + 0.2
                add_text_box(
                    q_slide,
                    option,
                    text_x,
                    oy,
                    option_width - (text_x - ox) - 0.2,
                    option_height,
                    font_size=18,
                    color=COLORS["text"],
                    alignment=PP_ALIGN.CENTER,
                    vertical_alignment=MSO_ANCHOR.MIDDLE,
                )
            presentation_title = content.get("title") or t("untitledPresentation")
            slide_number = slide_count_offset + quiz_slide_count
            add_footer(
                q_slide,
                presentation_title,
                slide_number,
                total_slides,
                THEME["footer_style"],
            )
            # Answer slide
            a_slide = prs.slides.add_slide(prs.slide_layouts[6])
            slides.append(a_slide)
            quiz_slide_count += 1
            add_shape(
                a_slide,
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                SLIDE_WIDTH,
                1.2,
                fill_color=COLORS["primary"],
            )
            add_text_box(
                a_slide,
                t("quizAnswer", num=q_idx + 1),
                0.5,
                0.3,
                9.0,
                0.6,
                font_size=40,
                bold=True,
                color=COLORS["text_light"],
                alignment=PP_ALIGN.CENTER,
                vertical_alignment=MSO_ANCHOR.MIDDLE,
            )
            add_shape(
                a_slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                0.5,
                1.4,
                9.0,
                0.8,
                fill_color=COLORS["light"],
                line_color=COLORS["light"],
                line_width=1,
            )
            add_text_box(
                a_slide,
                f"{t('question', text=question_text)}",
                0.7,
                1.5,
                8.6,
                0.6,
                font_size=18,
                italic=True,
                color=COLORS["text"],
            )
            correct_answer = question.get("correctAnswer", "")
            if correct_answer:
                add_shape(
                    a_slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    0.5,
                    2.4,
                    9.0,
                    1.0,
                    fill_color=COLORS["dark_alt"],
                    line_color=COLORS["success"],
                    line_width=3,
                )
                add_text_box(
                    a_slide,
                    t("correctAnswer"),
                    0.7,
                    2.5,
                    8.6,
                    0.4,
                    font_size=20,
                    bold=True,
                    color=COLORS["warning"],
                )
                add_text_box(
                    a_slide,
                    correct_answer,
                    0.7,
                    2.9,
                    8.6,
                    0.4,
                    font_size=18,
                    color=COLORS["text_light"],
                )
            explanation = question.get("explanation", "")
            if explanation:
                add_shape(
                    a_slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    0.5,
                    3.6,
                    9.0,
                    1.4,
                    fill_color=COLORS["light"],
                    line_color=COLORS["primary_light"],
                    line_width=1,
                )
                add_shape(
                    a_slide,
                    MSO_SHAPE.RECTANGLE,
                    0.7,
                    3.7,
                    0.1,
                    1.2,
                    fill_color=COLORS["primary"],
                )
                add_text_box(
                    a_slide,
                    t("explanation"),
                    0.9,
                    3.7,
                    8.5,
                    0.4,
                    font_size=20,
                    bold=True,
                    color=COLORS["text"],
                )
                add_text_box(
                    a_slide,
                    (
                        explanation
                        if isinstance(explanation, str)
                        else json.dumps(explanation)
                    ),
                    0.9,
                    4.2,
                    8.5,
                    0.7,
                    font_size=16,
                    color=COLORS["text"],
                )
            presentation_title = content.get("title") or t("untitledPresentation")
            slide_number = slide_count_offset + quiz_slide_count
            add_footer(
                a_slide,
                presentation_title,
                slide_number,
                total_slides,
                THEME["footer_style"],
            )
    return slides
