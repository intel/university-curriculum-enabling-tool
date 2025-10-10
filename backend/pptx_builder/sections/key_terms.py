from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from ..constants import COLORS, THEME, FOOTER_Y, SLIDE_WIDTH
from ..shapes import (
    add_gradient_background,
    add_corner_accent,
    add_text_box,
    add_shape,
    add_table,
    add_footer,
)
from ..localization import t


def create_key_terms_slide(prs, content, total_slides):
    key_terms = content.get("keyTerms", [])
    if not key_terms:
        return []
    terms_per_slide = 4
    total_terms = len(key_terms)
    slides_needed = (total_terms + terms_per_slide - 1) // terms_per_slide
    slides = []
    for slide_idx in range(slides_needed):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slides.append(slide)
        add_gradient_background(
            prs, slide, COLORS["primary"], COLORS["primary_dark"], angle=0
        )
        add_corner_accent(slide, COLORS["accent3"], 1.0, "bottom-left")
        title = t("keyTerms")
        if slide_idx > 0:
            title += t("continued")
        add_text_box(
            slide,
            title,
            0.5,
            0.1,
            9,
            0.6,
            font_size=36,
            bold=True,
            color=COLORS["text_light"],
        )
        start_idx = slide_idx * terms_per_slide
        end_idx = min(start_idx + terms_per_slide, total_terms)
        terms_for_slide = key_terms[start_idx:end_idx]
        table_height = min(3.5, 0.8 * (len(terms_for_slide) + 1))
        if THEME["content_box_shadow"]:
            add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                0.4,
                0.9,
                9.2,
                table_height + 0.2,
                fill_color=COLORS["light_alt"],
                shadow=True,
            )
        table = add_table(
            slide,
            len(terms_for_slide) + 1,
            2,
            0.5,
            1.0,
            9.0,
            table_height,
            header_bg_color=COLORS["royal_blue"],
            alt_row_bg_color=COLORS["light"],
            border_color=COLORS["primary_light"],
        )
        headers = [
            {
                "text": t("term"),
                "options": {
                    "fill": {"color": COLORS["royal_blue"]},
                    "color": COLORS["text_light"],
                    "fontSize": 18,
                    "bold": True,
                    "align": "center",
                },
            },
            {
                "text": t("definition"),
                "options": {
                    "fill": {"color": COLORS["royal_blue"]},
                    "color": COLORS["text_light"],
                    "fontSize": 18,
                    "bold": True,
                    "align": "center",
                },
            },
        ]
        for j, cell_data in enumerate(headers):
            cell = table.cell(0, j)
            p = cell.text_frame.paragraphs[0]
            p.text = cell_data["text"]
            opts = cell_data["options"]
            if "fill" in opts and "color" in opts["fill"]:
                cell.fill.solid()
                cell.fill.fore_color.rgb = opts["fill"]["color"]
            if p.runs:
                run = p.runs[0]
                run.font.color.rgb = opts["color"]
                run.font.size = Pt(opts["fontSize"])
                run.font.bold = opts["bold"]
            if opts.get("align") == "center":
                p.alignment = 1  # PP_ALIGN.CENTER value
        for i, term in enumerate(terms_for_slide):
            row_idx = i + 1
            even = i % 2 == 0
            bg = COLORS["background"] if even else COLORS["light"]
            term_cell = table.cell(row_idx, 0)
            term_cell.text = term.get("term", "")
            term_cell.fill.solid()
            term_cell.fill.fore_color.rgb = bg
            if term_cell.text_frame.paragraphs[0].runs:
                r = term_cell.text_frame.paragraphs[0].runs[0]
                r.font.bold = True
                r.font.size = Pt(16)
                r.font.color.rgb = COLORS["primary_dark"]
            def_cell = table.cell(row_idx, 1)
            def_cell.text = term.get("definition", "")
            def_cell.fill.solid()
            def_cell.fill.fore_color.rgb = bg
            if def_cell.text_frame.paragraphs[0].runs:
                dr = def_cell.text_frame.paragraphs[0].runs[0]
                dr.font.size = Pt(14)
        add_footer(
            slide,
            (content.get("title") or t("untitledPresentation")),
            slide_idx + 2,
            total_slides,
            THEME["footer_style"],
        )
    return slides
