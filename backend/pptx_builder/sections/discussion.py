import json
from pptx.enum.shapes import MSO_SHAPE
from ..constants import COLORS, SLIDE_WIDTH, FOOTER_Y, THEME
from ..shapes import add_shape, add_text_box, add_footer
from ..localization import t
from ..utils import estimate_text_height


def create_discussion_slides(prs, content, total_slides):
    assessment_ideas = content.get("assessmentIdeas", [])
    slides = []
    slide_count_offset = (
        2
        + len(content.get("keyTerms", [])) // 4
        - 1
        + len([s for i, s in enumerate(content.get("slides", [])) if i != 1])
        + len(content.get("activities", [])) * 2
        + 1
    )
    quiz_count = 0
    for idea in content.get("assessmentIdeas", []):
        if "quiz" in idea.get("type", "").lower():
            quiz_count += (
                len([q for q in idea.get("exampleQuestions", []) if q.get("options")])
                * 2
            )
    discussion_slide_count = 0
    for idea in assessment_ideas:
        if "discussion" not in idea.get("type", "").lower():
            continue
        for q_idx, question in enumerate(idea.get("exampleQuestions", [])):
            question_text = question.get("question", "Example question")
            guidance = question.get("correctAnswer", "")
            q_slide = prs.slides.add_slide(prs.slide_layouts[6])
            slides.append(q_slide)
            discussion_slide_count += 1
            add_shape(
                q_slide,
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                SLIDE_WIDTH,
                0.8,
                fill_color=COLORS["primary"],
            )
            add_text_box(
                q_slide,
                t("discussionQuestion", num=q_idx + 1),
                0.5,
                0.1,
                9.0,
                0.6,
                font_size=32,
                bold=True,
                color=COLORS["text_light"],
            )
            add_shape(
                q_slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                0.5,
                1.1,
                9.0,
                0.8,
                fill_color=COLORS["light"],
                line_color=COLORS["light"],
                line_width=1,
            )
            add_text_box(
                q_slide,
                question_text,
                0.7,
                1.2,
                8.6,
                0.6,
                font_size=20,
                bold=True,
                color=COLORS["text"],
            )
            question_text_height = estimate_text_height(question_text, 20, 8.6)
            next_y = 1.2 + question_text_height + 1.2
            add_shape(
                q_slide,
                MSO_SHAPE.RECTANGLE,
                0.7,
                next_y,
                0.1,
                0.4,
                fill_color=COLORS["primary"],
            )
            add_text_box(
                q_slide,
                t("groupDiscussion"),
                0.9,
                next_y,
                8.5,
                0.4,
                font_size=20,
                bold=True,
                color=COLORS["primary"],
            )
            add_text_box(
                q_slide,
                t("groupInstruction"),
                0.9,
                next_y + 0.5,
                8.5,
                0.4,
                font_size=18,
                color=COLORS["text"],
            )
            presentation_title = content.get("title") or t("untitledPresentation")
            slide_number = slide_count_offset + quiz_count + discussion_slide_count
            add_footer(
                q_slide,
                presentation_title,
                slide_number,
                total_slides,
                THEME["footer_style"],
            )
            a_slide = prs.slides.add_slide(prs.slide_layouts[6])
            slides.append(a_slide)
            discussion_slide_count += 1
            add_shape(
                a_slide,
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                SLIDE_WIDTH,
                0.8,
                fill_color=COLORS["primary"],
            )
            add_text_box(
                a_slide,
                t("facilitatorGuidance", num=q_idx + 1),
                0.5,
                0.1,
                9.0,
                0.6,
                font_size=32,
                bold=True,
                color=COLORS["text_light"],
            )
            add_shape(
                a_slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                0.5,
                1.1,
                9.0,
                0.8,
                fill_color=COLORS["light"],
                line_color=COLORS["light"],
                line_width=1,
            )
            add_text_box(
                a_slide,
                f"{t('question', text=question_text)}",
                0.7,
                1.2,
                8.6,
                0.6,
                font_size=18,
                italic=True,
                color=COLORS["text"],
            )
            guidance_y = 1.2 + question_text_height + 0.7
            if guidance:
                add_shape(
                    a_slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    0.5,
                    guidance_y,
                    9.0,
                    2.8,
                    fill_color=COLORS["light"],
                    line_color=COLORS["accent2"],
                    line_width=2,
                )
                add_shape(
                    a_slide,
                    MSO_SHAPE.RECTANGLE,
                    0.7,
                    guidance_y + 0.1,
                    0.1,
                    2.5,
                    fill_color=COLORS["accent2"],
                )
                add_text_box(
                    a_slide,
                    t("facilitatorGuidance", num=q_idx + 1).split(":")[0] + ":",
                    0.9,
                    guidance_y + 0.1,
                    8.5,
                    0.4,
                    font_size=20,
                    bold=True,
                    color=COLORS["accent2"],
                )
                add_text_box(
                    a_slide,
                    guidance if isinstance(guidance, str) else json.dumps(guidance),
                    0.9,
                    guidance_y + 0.6,
                    8.3,
                    1.5,
                    font_size=16,
                    color=COLORS["text"],
                )
            slide_number = slide_count_offset + quiz_count + discussion_slide_count
            add_footer(
                a_slide,
                presentation_title,
                slide_number,
                total_slides,
                THEME["footer_style"],
            )
    return slides
