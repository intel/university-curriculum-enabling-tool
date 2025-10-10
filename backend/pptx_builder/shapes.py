from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from .constants import COLORS, THEME, SLIDE_WIDTH, SLIDE_HEIGHT, FOOTER_Y


def add_text_box(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size=12,
    bold=False,
    italic=False,
    color=COLORS["text"],
    alignment=PP_ALIGN.LEFT,
    vertical_alignment=MSO_ANCHOR.TOP,
    level=0,
    bg_color=None,
    border_color=None,
    shadow=False,
):
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical_alignment
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except:
        pass
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.level = level
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = color
    if bg_color:
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    if border_color:
        line = textbox.line
        line.color.rgb = border_color
        line.width = Pt(1)
    if shadow and THEME["content_box_shadow"]:
        try:
            sh = textbox.shadow
            sh.inherit = False
            sh.visible = True
            sh.blur_radius = Pt(5)
            sh.distance = Pt(3)
            sh.angle = 45
            sh.color.rgb = RGBColor(0, 0, 0)
            sh.transparency = 0.7
        except:
            pass
    return textbox


def add_shape(
    slide,
    shape_type,
    left,
    top,
    width,
    height,
    fill_color=None,
    line_color=None,
    line_width=None,
    shadow=False,
    opacity=1.0,
):
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if opacity < 1.0:
            try:
                shape.fill.transparency = 1.0 - opacity
            except:
                pass
    if line_color:
        shape.line.color.rgb = line_color
    if line_width is not None:
        shape.line.width = Pt(line_width)
    if shadow:
        try:
            sh = shape.shadow
            sh.inherit = False
            sh.visible = True
            sh.blur_radius = Pt(5)
            sh.distance = Pt(3)
            sh.angle = 45
            sh.color.rgb = RGBColor(0, 0, 0)
            sh.transparency = 0.7
        except:
            pass
    return shape


def add_gradient_background(prs, slide, start_color, end_color, angle=90):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.line.fill.background()
    try:
        fill = shape.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = start_color
        fill.gradient_stops[0].position = 0
        fill.gradient_stops[1].color.rgb = end_color
        fill.gradient_stops[1].position = 1
        fill.gradient_angle = angle
    except:
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = start_color
    return shape


def add_corner_accent(slide, color=COLORS["accent1"], size=1.0, position="top-right"):
    if position == "top-right":
        left, top = SLIDE_WIDTH - size, 0
    elif position == "top-left":
        left, top = 0, 0
    elif position == "bottom-right":
        left, top = SLIDE_WIDTH - size, SLIDE_HEIGHT - size
    else:
        left, top = 0, SLIDE_HEIGHT - size
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    try:
        shape.fill.transparency = 0.3
    except:
        pass
    return shape


def add_table(slide, rows, cols, left, top, width, height, **kwargs):
    table = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    ).table
    return table


def add_footer(slide, title_text, slide_number, total_slides, style="modern"):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    if style == "modern":
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            0.5,
            FOOTER_Y - 0.05,
            9.0,
            0.01,
            fill_color=COLORS["primary_light"],
            opacity=0.5,
        )
        add_text_box(
            slide,
            title_text,
            0.5,
            FOOTER_Y,
            8.5,
            0.3,
            font_size=10,
            color=COLORS["primary"],
            italic=True,
        )
        add_text_box(
            slide,
            f"{slide_number}",
            9.0,
            FOOTER_Y,
            0.5,
            0.3,
            font_size=10,
            color=COLORS["primary"],
            alignment=PP_ALIGN.RIGHT,
        )
    else:
        add_text_box(
            slide,
            title_text,
            0.5,
            FOOTER_Y,
            8.5,
            0.3,
            font_size=10,
            color=COLORS["royal_blue"],
        )
        add_text_box(
            slide,
            f"{slide_number}",
            9.0,
            FOOTER_Y,
            0.5,
            0.3,
            font_size=10,
            color=COLORS["text"],
            alignment=PP_ALIGN.RIGHT,
        )
